from __future__ import annotations

from pathlib import Path
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "generated"
ASSET_DIR = OUT_DIR / "aesthetic_assets"
OUT_PPTX = OUT_DIR / "Saheli_WitchHunt_2026_Aesthetic_Deck.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_HEAD = "Noto Serif"
FONT_BODY = "Noto Sans"
FONT_DEV = "Noto Sans Devanagari"

PIL_HEAD_BOLD = "/usr/share/fonts/noto/NotoSerif-Bold.ttf"
PIL_BODY_BOLD = "/usr/share/fonts/noto/NotoSans-Bold.ttf"
PIL_BODY = "/usr/share/fonts/noto/NotoSans-Regular.ttf"

INK = RGBColor(15, 18, 25)
NAVY = RGBColor(18, 29, 43)
FOREST = RGBColor(45, 76, 66)
MOSS = RGBColor(113, 139, 113)
GOLD = RGBColor(208, 167, 96)
PLUM = RGBColor(102, 83, 120)
PAPER = RGBColor(244, 238, 228)
PAPER_2 = RGBColor(236, 229, 216)
WHITE = RGBColor(250, 248, 244)
MUTED_LIGHT = RGBColor(205, 204, 196)
MUTED_DARK = RGBColor(88, 92, 98)
RED = RGBColor(168, 78, 68)


def rgb(hex_str: str) -> tuple[int, int, int]:
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i : i + 2], 16) for i in (0, 2, 4))


def build_backgrounds():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    make_dark_bg(ASSET_DIR / "bg_dark.png")
    make_paper_bg(ASSET_DIR / "bg_paper.png")
    make_grid_bg(ASSET_DIR / "bg_grid.png")
    make_closing_bg(ASSET_DIR / "bg_closing.png")


def add_grain(img: Image.Image, alpha: int):
    px = img.load()
    w, h = img.size
    for _ in range(4200):
        x = random.randint(0, w - 1)
        y = random.randint(0, h - 1)
        val = random.randint(190, 255)
        px[x, y] = (val, val, val, alpha)


def make_dark_bg(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#0d1117"))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.ellipse((1040, -180, 1700, 420), fill=(208, 167, 96, 48))
    d.ellipse((-220, 520, 420, 1080), fill=(91, 123, 105, 40))
    d.ellipse((1180, 620, 1500, 930), fill=(101, 83, 122, 36))
    d.arc((930, -60, 1550, 560), 18, 340, fill=(214, 178, 111, 185), width=5)
    d.arc((960, -30, 1520, 530), 18, 340, fill=(119, 143, 117, 135), width=2)
    d.arc((1160, 600, 1520, 960), 15, 310, fill=(100, 84, 120, 160), width=4)
    d.line((1160, 70, 1450, 70), fill=(255, 255, 255, 18), width=1)
    d.line((120, 760, 430, 760), fill=(255, 255, 255, 20), width=1)
    overlay = overlay.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, overlay)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_grain(grain, 12)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_paper_bg(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#f5efe4"))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.ellipse((1050, -120, 1700, 320), fill=(213, 184, 132, 38))
    d.ellipse((-260, 640, 360, 1060), fill=(123, 151, 126, 34))
    d.arc((1120, 20, 1550, 450), 12, 310, fill=(90, 118, 104, 90), width=3)
    d.arc((1170, 70, 1500, 400), 12, 310, fill=(213, 170, 97, 120), width=2)
    overlay = overlay.filter(ImageFilter.GaussianBlur(12))
    img = Image.alpha_composite(img, overlay)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_grain(grain, 8)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_grid_bg(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#101827"))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    for x in range(0, 1600, 70):
        d.line((x, 0, x, 900), fill=(90, 112, 144, 25), width=1)
    for y in range(0, 900, 70):
        d.line((0, y, 1600, y), fill=(90, 112, 144, 25), width=1)
    d.ellipse((960, 120, 1520, 680), fill=(55, 98, 123, 48))
    d.ellipse((-180, 500, 380, 980), fill=(95, 76, 120, 32))
    d.rounded_rectangle((890, 230, 1500, 600), radius=48, outline=(210, 169, 100, 90), width=3)
    overlay = overlay.filter(ImageFilter.GaussianBlur(14))
    img = Image.alpha_composite(img, overlay)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_grain(grain, 10)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_closing_bg(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#0d1117"))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.arc((980, -80, 1660, 600), 20, 340, fill=(212, 169, 95, 205), width=6)
    d.arc((1010, -40, 1630, 570), 20, 340, fill=(102, 83, 120, 135), width=3)
    d.arc((1080, 560, 1580, 1060), 10, 300, fill=(92, 124, 106, 150), width=4)
    d.ellipse((1100, 40, 1480, 420), fill=(205, 169, 101, 38))
    d.ellipse((1200, 620, 1600, 980), fill=(96, 76, 121, 28))
    overlay = overlay.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, overlay)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_grain(grain, 12)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide_obj, image_name: str):
    slide_obj.shapes.add_picture(
        str(ASSET_DIR / image_name),
        0,
        0,
        width=Inches(SLIDE_W),
        height=Inches(SLIDE_H),
    )


def set_run_style(run, font_name: str, size: float, color: RGBColor, bold: bool = False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide_obj,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    font_name: str = FONT_BODY,
    size: float = 18,
    color: RGBColor = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.05,
):
    box = slide_obj.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_run_style(r, font_name, size, color, bold)
    return box


def add_rich_lines(
    slide_obj,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, dict]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    box = slide_obj.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.line_spacing = opts.get("line_spacing", 1.02)
        r = p.add_run()
        r.text = text
        set_run_style(
            r,
            opts.get("font_name", FONT_BODY),
            opts.get("size", 18),
            opts.get("color", WHITE),
            opts.get("bold", False),
        )
    return box


def add_bullets(
    slide_obj,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[str],
    *,
    font_name: str = FONT_BODY,
    size: float = 15,
    color: RGBColor = WHITE,
    bullet_color: RGBColor = GOLD,
    spacing: float = 1.18,
):
    box = slide_obj.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.line_spacing = spacing
        r1 = p.add_run()
        r1.text = "• "
        set_run_style(r1, font_name, size, bullet_color, True)
        r2 = p.add_run()
        r2.text = item
        set_run_style(r2, font_name, size, color, False)
    return box


def add_card(
    slide_obj,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius_shape: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    shape = slide_obj.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    else:
        shape.line.color.rgb = fill
        shape.line.width = Pt(0)
    return shape


def add_chip(slide_obj, x, y, w, h, text, *, fill, text_color, line=None, size=11):
    add_card(slide_obj, x, y, w, h, fill=fill, line=line or fill)
    add_text(slide_obj, x, y + 0.05, w, h - 0.05, text, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def add_rule(slide_obj, x, y, w, color):
    line = slide_obj.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    return line


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def build_slide_1(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_dark.png")
    add_text(s, 0.75, 0.55, 2.3, 0.3, "WITCHHUNT 2026 PROPOSAL", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.75, 1.05, 5.2, 1.15, "SAHELI", font_name=FONT_HEAD, size=32, color=WHITE, bold=True)
    add_text(s, 0.78, 2.03, 2.2, 0.42, "सहेली", font_name=FONT_DEV, size=21, color=RGBColor(230, 221, 204), bold=True)
    add_text(
        s,
        0.75,
        2.48,
        5.6,
        1.15,
        "Voice-native, stealth-first rights and evidence infrastructure\nfor women on shared phones in India.",
        size=18,
        color=WHITE,
        bold=False,
    )
    add_card(s, 0.75, 4.0, 4.8, 1.35, fill=RGBColor(33, 41, 52), line=RGBColor(89, 98, 112))
    add_text(
        s,
        1.0,
        4.28,
        4.3,
        0.72,
        "Not a chatbot. Not a wrapper.\nA privacy-first access layer for rights, schemes, health, and emergency action.",
        size=16,
        color=WHITE,
    )
    chips = [
        ("Evidence Vault", 0.75, 6.25, 1.65),
        ("Rights Engine", 2.55, 6.25, 1.6),
        ("Scheme Navigator", 4.28, 6.25, 1.95),
        ("Health Journal", 6.38, 6.25, 1.7),
        ("Stealth SOS", 8.25, 6.25, 1.45),
    ]
    for label, x, y, w in chips:
        add_chip(s, x, y, w, 0.46, label, fill=RGBColor(30, 35, 44), text_color=WHITE, line=RGBColor(74, 82, 95), size=10)
    add_text(s, 9.65, 5.9, 2.5, 0.65, "HopeWorks x AI4India x Skills Cafe", size=11, color=MUTED_LIGHT, align=PP_ALIGN.RIGHT)


def build_slide_2(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 2.7, 0.32, "THE ACCESS CRISIS", size=11, color=MUTED_DARK, bold=True)
    add_text(s, 0.8, 0.95, 5.6, 0.9, "Women’s empowerment is an information-access problem trapped inside a privacy crisis.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    stats = [
        ("43%", "Indian women without a personal smartphone"),
        ("58%", "Rural women using shared devices"),
        ("67%", "Harassment that still goes unreported"),
        ("26%", "DV conviction rate constrained by proof gaps"),
        ("74%", "Nirbhaya Fund still left unspent"),
    ]
    positions = [(0.85, 2.2), (2.95, 2.2), (5.05, 2.2), (0.85, 4.35), (2.95, 4.35)]
    for (num, label), (x, y) in zip(stats, positions, strict=True):
        add_card(s, x, y, 1.8, 1.5, fill=WHITE, line=RGBColor(214, 203, 186))
        add_text(s, x + 0.12, y + 0.18, 1.55, 0.45, num, font_name=FONT_HEAD, size=24, color=FOREST, bold=True)
        add_text(s, x + 0.12, y + 0.75, 1.5, 0.55, label, size=9.5, color=MUTED_DARK)
    add_card(s, 7.3, 1.8, 5.2, 4.55, fill=RGBColor(27, 39, 52), line=RGBColor(67, 84, 102))
    add_text(s, 7.65, 2.1, 4.5, 0.28, "CORE INSIGHT", size=11, color=RGBColor(218, 201, 169), bold=True)
    add_text(
        s,
        7.65,
        2.45,
        4.45,
        1.8,
        "The laws exist. The schemes exist. The support infrastructure exists.\n\n"
        "What does not exist is a safe way for women on shared phones to access all of it privately, in plain language, and without reliable internet.",
        size=17,
        color=WHITE,
    )
    add_bullets(
        s,
        7.65,
        4.65,
        4.4,
        1.35,
        [
            "Shared devices turn visible safety apps into liabilities.",
            "Complex English-heavy UI excludes WhatsApp-first users.",
            "Emergency-only tools ignore proof, rights, and recovery.",
        ],
        size=12,
        color=RGBColor(236, 234, 227),
        bullet_color=GOLD,
    )
    add_text(s, 0.8, 6.9, 4.8, 0.2, "Sources: NSSO, IAMAI, GSMA, NCRB, NARI, NALSA", size=9.5, color=MUTED_DARK)


def build_slide_3(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_dark.png")
    add_text(s, 0.8, 0.55, 3.0, 0.3, "WHY CURRENT TOOLS FAIL", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 0.95, 5.4, 0.8, "Most women’s safety products break exactly where privacy and access matter most.", font_name=FONT_HEAD, size=25, color=WHITE, bold=True)
    cards = [
        ("01", "Assume private phone ownership", "A visible safety app can itself become dangerous on a monitored device."),
        ("02", "Depend on continuous internet", "GPS-heavy and chat-heavy flows break in low-connectivity conditions."),
        ("03", "React only after escalation", "SOS without evidence, rights, or recovery support leaves the core problem unsolved."),
        ("04", "Route women into weak response loops", "Trust in formal reporting remains low, and support pathways are fragmented."),
        ("05", "Assume digital literacy", "If it feels harder than a voice note, many users will never reach the useful part."),
    ]
    positions = [(0.85, 2.1), (4.35, 2.1), (7.85, 2.1), (2.6, 4.5), (6.1, 4.5)]
    for (num, title, body), (x, y) in zip(cards, positions, strict=True):
        add_card(s, x, y, 2.95, 1.65, fill=RGBColor(27, 34, 44), line=RGBColor(75, 86, 99))
        add_text(s, x + 0.18, y + 0.18, 0.5, 0.35, num, size=11, color=GOLD, bold=True)
        add_text(s, x + 0.18, y + 0.48, 2.5, 0.5, title, size=13.5, color=WHITE, bold=True)
        add_text(s, x + 0.18, y + 0.95, 2.5, 0.5, body, size=10.5, color=MUTED_LIGHT)
    add_card(s, 0.85, 6.45, 11.6, 0.52, fill=RGBColor(31, 40, 49), line=RGBColor(83, 91, 102))
    add_text(s, 1.1, 6.55, 11.1, 0.22, "Saheli’s bet: if support begins with conversation, camouflage, and offline tolerance, excluded women can finally act.", size=11.5, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_4(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 2.7, 0.3, "WHAT SAHELI IS", size=11, color=MUTED_DARK, bold=True)
    add_text(s, 0.8, 0.95, 4.9, 0.9, "One sentence in. A structured support journey out.", font_name=FONT_HEAD, size=26, color=INK, bold=True)
    add_card(s, 0.8, 1.9, 4.8, 4.8, fill=RGBColor(25, 37, 50), line=RGBColor(70, 89, 108))
    add_text(s, 1.1, 2.18, 4.2, 0.24, "ONE-LINE PITCH", size=11, color=RGBColor(218, 201, 169), bold=True)
    add_text(
        s,
        1.1,
        2.55,
        4.1,
        1.75,
        "Saheli is a voice-native AI companion hidden inside a devotional cover app.\n\n"
        "It gives women access to an encrypted evidence vault, step-by-step rights guidance, scheme eligibility, health screening, and emergency SOS in Hindi, offline, on shared phones.",
        size=17,
        color=WHITE,
    )
    add_text(s, 1.1, 5.55, 4.1, 0.55, "Built for women 18-45 in tier-2/3 India, including shared-phone and feature-phone users.", size=11.5, color=MUTED_LIGHT)
    pillars = [
        ("Voice-native", "Conversation is the primary interface, not menus."),
        ("Stealth-first", "Hidden behind a devotional app, with quick exit and covert entry."),
        ("Offline-first", "SMS, cached PWA, and low-end device paths are built in."),
        ("Evidence-first", "Proof capture is encrypted, hashed, and usable later."),
    ]
    coords = [(6.2, 1.95), (9.35, 1.95), (6.2, 4.35), (9.35, 4.35)]
    for (title, body), (x, y) in zip(pillars, coords, strict=True):
        add_card(s, x, y, 2.75, 1.85, fill=WHITE, line=RGBColor(214, 203, 186))
        add_text(s, x + 0.18, y + 0.22, 2.3, 0.36, title, size=13.5, color=FOREST, bold=True)
        add_text(s, x + 0.18, y + 0.7, 2.3, 0.72, body, size=10.5, color=MUTED_DARK)


def build_slide_5(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_dark.png")
    add_text(s, 0.8, 0.55, 3.4, 0.3, "ONE CONVERSATION, FIVE CAPABILITY LAYERS", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 0.95, 4.8, 0.8, "Users should not have to learn a system before they can use it.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    add_text(s, 0.8, 1.9, 3.6, 1.0, "Saheli uses the familiarity of a voice note, then routes the user into the right workflow beneath the surface.", size=15.5, color=MUTED_LIGHT)
    items = [
        ("Evidence Vault", "“Saboot rakhna hai”", "Photos, threats, documents, timestamps, GPS, local encryption."),
        ("Know Your Rights", "“Mera kya haq hai?”", "PWDVA, POSH, maintenance, property, custody, legal aid flows."),
        ("Safety Net", "“Mujhe madad chahiye”", "Stealth SOS, helplines, OSC locator, offline SMS fallback."),
        ("Entitlements", "“Sarkari yojana”", "Eligibility checks for welfare schemes, loans, and support funds."),
        ("Health Journal", "“Tabiyat theek nahi”", "Voice screening, symptom tracking, and doctor-ready summaries."),
    ]
    y = 1.65
    for title, trigger, body in items:
        add_card(s, 5.4, y, 6.85, 0.9, fill=RGBColor(28, 35, 44), line=RGBColor(76, 86, 98))
        add_text(s, 5.72, y + 0.13, 2.0, 0.22, title, size=13.5, color=WHITE, bold=True)
        add_text(s, 7.82, y + 0.14, 1.45, 0.2, trigger, size=11, color=GOLD, bold=True)
        add_text(s, 5.72, y + 0.44, 5.9, 0.22, body, size=10.2, color=MUTED_LIGHT)
        y += 0.98


def build_slide_6(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_grid.png")
    add_text(s, 0.8, 0.55, 2.7, 0.3, "STEALTH EXPERIENCE FLOW", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 0.95, 5.1, 0.8, "The product disappears when it must, and becomes actionable when it counts.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    steps = [
        ("01", "Install cover app", "Looks like Daily Mantras, not a safety tool."),
        ("02", "Hidden unlock", "Five taps plus PIN opens the real interface."),
        ("03", "Speak naturally", "Hindi-first voice queries replace navigation."),
        ("04", "Route to action", "Rights, evidence, schemes, or SOS workflows load."),
        ("05", "Exit cleanly", "Quick exit and offline fallback reduce visible traces."),
    ]
    x_positions = [0.95, 3.35, 5.75, 8.15, 10.55]
    for idx, ((num, title, body), x) in enumerate(zip(steps, x_positions, strict=True)):
        add_card(s, x, 3.1, 2.0, 2.35, fill=RGBColor(23, 31, 43), line=RGBColor(78, 91, 110))
        add_text(s, x + 0.18, 3.28, 0.48, 0.3, num, size=11.5, color=GOLD, bold=True)
        add_text(s, x + 0.18, 3.62, 1.55, 0.45, title, size=13, color=WHITE, bold=True)
        add_text(s, x + 0.18, 4.18, 1.55, 0.72, body, size=10.2, color=MUTED_LIGHT)
        if idx < len(steps) - 1:
            add_rule(s, x + 1.9, 4.22, 0.48, RGBColor(97, 108, 124))
    add_card(s, 0.95, 6.25, 11.6, 0.48, fill=RGBColor(27, 36, 47), line=RGBColor(76, 87, 98))
    add_text(s, 1.15, 6.34, 11.2, 0.18, "No visible safety app. No complex navigation. No assumption of continuous internet.", size=11.5, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_7(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 2.0, 0.3, "THREE USERS", size=11, color=MUTED_DARK, bold=True)
    add_text(s, 0.8, 0.95, 5.8, 0.8, "Different entry points. The same underlying support system.", font_name=FONT_HEAD, size=25, color=INK, bold=True)
    personas = [
        ("Priya", "Shared-phone DV survivor", ["Entry: health conversation under cover mode", "Need: safe proof + rights + covert help", "Outcome: evidence, maintenance guidance, OSC path"], PLUM),
        ("Sunita", "Teacher with severe period pain", ["Entry: symptom description in plain Hindi", "Need: screening + low-cost care + doctor card", "Outcome: structured tracking and affordable treatment path"], FOREST),
        ("Meena", "Feature-phone user without internet", ["Entry: SMS keyword and guided prompts", "Need: no-install screening and referral", "Outcome: PHC direction and ASHA follow-up"], RGBColor(127, 124, 79)),
    ]
    xs = [0.85, 4.45, 8.05]
    for (name, subtitle, points, accent), x in zip(personas, xs, strict=True):
        add_card(s, x, 2.1, 3.0, 4.55, fill=WHITE, line=RGBColor(214, 203, 186))
        add_card(s, x, 2.1, 3.0, 0.18, fill=accent, line=accent, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        add_text(s, x + 0.22, 2.38, 2.4, 0.4, name, font_name=FONT_HEAD, size=20, color=INK, bold=True)
        add_text(s, x + 0.22, 2.82, 2.4, 0.3, subtitle, size=11, color=MUTED_DARK, bold=True)
        add_bullets(s, x + 0.22, 3.35, 2.5, 2.1, points, size=10.6, color=MUTED_DARK, bullet_color=accent, spacing=1.16)
        add_chip(s, x + 0.22, 5.88, 1.15, 0.38, "USER", fill=PAPER_2, text_color=MUTED_DARK, line=RGBColor(220, 210, 194), size=9)
        add_chip(s, x + 1.5, 5.88, 1.2, 0.38, "REAL NEED", fill=PAPER_2, text_color=MUTED_DARK, line=RGBColor(220, 210, 194), size=9)


def build_slide_8(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_grid.png")
    add_text(s, 0.8, 0.55, 2.9, 0.3, "SOVEREIGN AI + PUBLIC INFRASTRUCTURE", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 0.95, 6.1, 0.8, "Saheli works because the model layer and the delivery layer are both localized.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    add_card(s, 0.85, 1.95, 5.85, 4.95, fill=RGBColor(19, 28, 39), line=RGBColor(73, 86, 104))
    add_text(s, 1.15, 2.2, 2.5, 0.25, "AI STACK", size=11, color=GOLD, bold=True)
    chips_left = [
        "Sarvam Saaras V3  |  Indic STT",
        "Bulbul V3  |  Voice output",
        "Sarvam-30B + legal RAG",
        "Sarvam Edge  |  offline AI",
        "Sarvam Vision  |  OCR",
        "Samvaad  |  WhatsApp, web, SMS",
        "React Native + PWA",
        "SHA-256 evidence pipeline",
    ]
    y = 2.65
    for label in chips_left:
        add_chip(s, 1.15, y, 4.95, 0.42, label, fill=RGBColor(28, 38, 50), text_color=WHITE, line=RGBColor(80, 92, 111), size=10)
        y += 0.54
    add_card(s, 6.95, 1.95, 5.5, 4.95, fill=RGBColor(23, 32, 42), line=RGBColor(73, 86, 104))
    add_text(s, 7.25, 2.2, 3.7, 0.25, "DISTRIBUTION + SUPPORT LAYER", size=11, color=GOLD, bold=True)
    chips_right = [
        "733+ One Stop Centres",
        "9M Self-Help Groups",
        "1M+ ASHA workers",
        "District legal aid networks",
        "9,000+ Jan Aushadhi stores",
        "ABHA-linked digital health records",
    ]
    y = 2.65
    for label in chips_right:
        add_chip(s, 7.25, y, 4.55, 0.42, label, fill=RGBColor(28, 38, 50), text_color=WHITE, line=RGBColor(80, 92, 111), size=10)
        y += 0.54
    add_text(s, 7.25, 6.0, 4.5, 0.44, "The app is not the intervention by itself; it is the coordination layer across existing institutions.", size=11, color=MUTED_LIGHT)


def build_slide_9(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 2.4, 0.3, "COMPETITIVE EDGE", size=11, color=MUTED_DARK, bold=True)
    add_text(s, 0.8, 0.95, 5.0, 0.8, "Saheli is not another generic safety app.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    columns = ["Feature", "112 India", "Himmat", "VictimsVoice", "Saheli"]
    rows = [
        ("Stealth mode", "No", "No", "No", "Yes"),
        ("Offline path", "No", "No", "No", "Yes"),
        ("Voice-native UX", "No", "No", "No", "Yes"),
        ("Evidence vault", "No", "No", "Yes", "Yes"),
        ("Rights workflows", "No", "No", "No", "Yes"),
        ("Scheme navigation", "No", "No", "No", "Yes"),
    ]
    widths = [2.55, 1.65, 1.45, 1.85, 1.75]
    x0 = 0.95
    y0 = 2.0
    x = x0
    for col, w in zip(columns, widths, strict=True):
        add_card(s, x, y0, w, 0.58, fill=INK if col == "Feature" else FOREST, line=INK if col == "Feature" else FOREST, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        add_text(s, x, y0 + 0.14, w, 0.22, col, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w + 0.08
    y = y0 + 0.72
    for feature, c1, c2, c3, c4 in rows:
        x = x0
        values = [feature, c1, c2, c3, c4]
        for idx, (val, w) in enumerate(zip(values, widths, strict=True)):
            fill = WHITE if idx == 0 else PAPER_2
            line = RGBColor(217, 208, 193)
            if idx == 4:
                fill = RGBColor(225, 235, 227)
                line = RGBColor(171, 196, 176)
            add_card(s, x, y, w, 0.52, fill=fill, line=line, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            color = INK if idx == 0 else (FOREST if val == "Yes" else MUTED_DARK)
            bold = idx == 0 or val == "Yes"
            add_text(s, x, y + 0.14, w, 0.18, val, size=10.5, color=color, bold=bold, align=PP_ALIGN.CENTER)
            x += w + 0.08
        y += 0.58
    add_card(s, 0.95, 6.25, 11.25, 0.48, fill=INK, line=INK)
    add_text(s, 1.2, 6.34, 10.8, 0.18, "Saheli is the only concept here designed for shared phones, low digital literacy, and offline access at the same time.", size=11.2, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_10(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_dark.png")
    add_text(s, 0.8, 0.55, 2.3, 0.3, "YEAR 1 IMPACT", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 0.95, 4.6, 0.8, "Measure the platform by documented action, not downloads.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    cards = [
        ("10,000+", "women using the evidence vault"),
        ("50,000+", "rights journeys completed"),
        ("25,000+", "eligibility checks for schemes"),
        ("100,000+", "voice-based health screenings"),
    ]
    xs = [0.95, 3.95, 6.95, 9.95]
    for (num, label), x in zip(cards, xs, strict=True):
        add_card(s, x, 2.15, 2.35, 1.62, fill=RGBColor(28, 35, 44), line=RGBColor(79, 90, 103))
        add_text(s, x + 0.14, 2.36, 2.05, 0.42, num, font_name=FONT_HEAD, size=23, color=GOLD, bold=True)
        add_text(s, x + 0.14, 2.95, 2.0, 0.42, label, size=10.5, color=MUTED_LIGHT)
    add_card(s, 0.95, 4.25, 5.6, 2.05, fill=RGBColor(23, 31, 42), line=RGBColor(77, 89, 104))
    add_text(s, 1.25, 4.5, 2.4, 0.25, "SUSTAINABILITY CHANNELS", size=11, color=GOLD, bold=True)
    add_bullets(
        s,
        1.25,
        4.85,
        4.9,
        1.1,
        [
            "ASHA workers introduce the tool in routine visits.",
            "SHG leaders run small-group onboarding sessions.",
            "OSCs and legal aid teams use Saheli as a case-support layer.",
        ],
        size=11.2,
        color=MUTED_LIGHT,
        bullet_color=GOLD,
    )
    add_card(s, 6.85, 4.25, 5.4, 2.05, fill=RGBColor(23, 31, 42), line=RGBColor(77, 89, 104))
    add_text(s, 7.15, 4.5, 2.4, 0.25, "REVENUE WITHOUT PAYWALLING THE USER", size=11, color=GOLD, bold=True)
    add_bullets(
        s,
        7.15,
        4.85,
        4.7,
        1.1,
        [
            "B2G tools for ASHA and OSC workflows.",
            "B2B POSH compliance and workplace support layers.",
            "Grant-backed scale-up while keeping the core user product free.",
        ],
        size=11.2,
        color=MUTED_LIGHT,
        bullet_color=GOLD,
    )
    add_text(s, 0.95, 6.75, 4.2, 0.22, "Feature phone target: 20,000+ SMS interactions", size=10.5, color=MUTED_LIGHT)


def build_slide_11(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 3.4, 0.3, "FACILITATED DATA REQUEST", size=11, color=MUTED_DARK, bold=True)
    add_text(s, 0.8, 0.95, 5.6, 0.8, "What would materially improve the prototype during the challenge.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    reqs = [
        ("Indic speech from women in low-resource dialects", "Needed to validate STT quality for rural Hindi and code-mixed usage."),
        ("Verified district-level service directories", "Needed for trustworthy routing to OSCs, PHCs, legal aid, and Jan Aushadhi."),
        ("Anonymized women’s health validation datasets", "Needed to calibrate screening flows and reduce false reassurance."),
    ]
    xs = [0.95, 4.32, 7.69]
    for (title, body), x in zip(reqs, xs, strict=True):
        add_card(s, x, 2.05, 2.95, 3.65, fill=WHITE, line=RGBColor(214, 203, 186))
        add_text(s, x + 0.18, 2.28, 2.5, 0.7, title, size=13.2, color=FOREST, bold=True)
        add_text(s, x + 0.18, 3.15, 2.45, 1.0, body, size=10.8, color=MUTED_DARK)
    add_card(s, 0.95, 6.0, 10.95, 0.55, fill=INK, line=INK)
    add_text(s, 1.25, 6.12, 10.45, 0.22, "Preferred delivery: secure API, sandbox, or encrypted batch access with DPDP-aligned handling.", size=11.3, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_12(prs: Presentation):
    s = slide(prs)
    add_bg(s, "bg_closing.png")
    add_text(s, 0.8, 0.75, 4.5, 0.32, "WHY NOW", size=11, color=MUTED_LIGHT, bold=True)
    add_text(s, 0.8, 1.2, 6.0, 1.2, "Build for the women current safety apps were never designed for.", font_name=FONT_HEAD, size=28, color=WHITE, bold=True)
    add_text(
        s,
        0.8,
        3.0,
        5.5,
        1.0,
        "Shared phones. Intermittent internet.\nHindi-first interaction.\nA real need for proof, rights, and help.",
        size=18,
        color=MUTED_LIGHT,
    )
    add_card(s, 0.8, 5.2, 5.0, 1.0, fill=RGBColor(28, 35, 44), line=RGBColor(79, 90, 103))
    add_text(s, 1.1, 5.48, 4.4, 0.38, "One conversation. Every right. Every scheme. Every tool.", size=18, color=WHITE, bold=True)
    add_text(s, 9.0, 6.35, 3.0, 0.42, "THANK YOU", font_name=FONT_HEAD, size=23, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)


def build():
    build_backgrounds()
    prs = make_presentation()
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))


if __name__ == "__main__":
    random.seed(7)
    build()
