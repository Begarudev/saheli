#!/usr/bin/env python3
"""
Saheli — WitchHunt 2026 Solution Proposal Deck  v3
Fixes applied:
  - Persona slides: panels now fill full 7.5" height (no dead space)
  - Title / AI-badge overlaps fixed on slides 5, 9, 11, 13
  - Slide 12: top strip reduced 2.9→1.6", challenges no longer cut off
  - Slide 2: removed purposeless blobs; stat cards taller & better spaced
  - Slide 11: closed 1.5" gap between body and tech-stack list
  - Slide 4/5 right panel: mandala fills more space, less empty filler
  - Consistent 1.0" left margin after AI badge on all left-title slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math, os

IMGDIR = "/home/garudev/dev/saheli/assets/images"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
MIDNIGHT   = RGBColor(0x0B, 0x07, 0x24)
INDIGO     = RGBColor(0x1A, 0x0A, 0x45)
MAROON     = RGBColor(0x5C, 0x0F, 0x26)
SAFFRON    = RGBColor(0xE8, 0x92, 0x0A)
GOLD_LIGHT = RGBColor(0xF5, 0xC5, 0x42)
LOTUS      = RGBColor(0xD4, 0x59, 0x7A)
TEAL       = RGBColor(0x00, 0x88, 0x88)
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)
ROSE_GOLD  = RGBColor(0xC9, 0x83, 0x6B)
DEEP_CARD  = RGBColor(0x12, 0x08, 0x32)

# ── Primitive helpers ─────────────────────────────────────────────────────────
def fill_bg(slide, r, g, b):
    bg = slide.background; f = bg.fill
    f.solid(); f.fore_color.rgb = RGBColor(r, g, b)

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=CREAM,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return txb

def img(slide, fname, l, t, w, h):
    """Insert a real photo/art image. fname is relative to IMGDIR."""
    path = os.path.join(IMGDIR, fname)
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def gold_line(slide, l, t, w, thick=0.03):
    rect(slide, l, t, w, thick, fill=SAFFRON)

def logo_bar(slide):
    txt(slide, "HopeWorks  ✦  AI4India",
        9.6, 0.1, 3.5, 0.4, size=9, color=ROSE_GOLD, align=PP_ALIGN.RIGHT)

def ai_gem(slide):
    """Small AI badge — top-left, occupies x 0.12-0.80, y 0.10-0.78."""
    oval(slide, 0.12, 0.1, 0.68, 0.68, fill=MAROON, line=SAFFRON, lw=Pt(1.5))
    txt(slide, "AI", 0.12, 0.1, 0.68, 0.68,
        size=15, bold=True, color=SAFFRON, align=PP_ALIGN.CENTER)

def gold_pill(slide, label, l, t, w=2.2):
    rect(slide, l, t, w, 0.38, fill=SAFFRON)
    txt(slide, label, l+0.08, t+0.03, w-0.12, 0.32,
        size=10, bold=True, color=MIDNIGHT)

# ── Decorative geometry ───────────────────────────────────────────────────────
def corner_mandala(slide, cx, cy, max_r=0.9, rings=4):
    colors = [SAFFRON, LOTUS, TEAL, ROSE_GOLD, INDIGO]
    for i in range(rings, 0, -1):
        r = max_r * i / rings
        oval(slide, cx-r/2, cy-r/2, r, r,
             line=colors[i % len(colors)], lw=Pt(1.5 if i < rings else 0.8))
    oval(slide, cx-0.07, cy-0.07, 0.14, 0.14, fill=SAFFRON)

def petal_ring(slide, cx, cy, r=0.6, n=8, color=LOTUS):
    pr = 0.18
    for i in range(n):
        angle = 2*math.pi*i/n
        oval(slide, cx+r*math.cos(angle)-pr/2, cy+r*math.sin(angle)-pr/2,
             pr, pr, line=color, lw=Pt(1))

def rich_mandala(slide, cx, cy, max_r=3.2, rings=6):
    """Larger, richer mandala for right-side decoration panels."""
    colors = [SAFFRON, LOTUS, TEAL, ROSE_GOLD, GOLD_LIGHT, INDIGO]
    for i in range(rings, 0, -1):
        r = max_r * i / rings
        oval(slide, cx-r/2, cy-r/2, r, r,
             line=colors[i % len(colors)],
             lw=Pt(1.8 if i == rings else (1.2 if i > rings//2 else 0.7)))
    # inner petals at 40% radius
    petal_ring(slide, cx, cy, r=max_r*0.4, n=12, color=LOTUS)
    petal_ring(slide, cx, cy, r=max_r*0.7, n=16, color=SAFFRON)
    oval(slide, cx-0.1, cy-0.1, 0.2, 0.2, fill=SAFFRON)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Warli painting fills right panel — real tribal art replaces decorative ovals
img(sl, "r_warli_jivya.jpg", 6.8, 0.0, 6.53, 7.5)
# Subtle vignette rings layered on top of photo
oval(sl, 6.5, -1.8, 8.5, 11.0, line=INDIGO, lw=Pt(1.5))
oval(sl, 7.5, -0.5, 6.5, 8.5,  line=MAROON, lw=Pt(1))

# Warm accent stripe left
rect(sl, 0.0, 0.0, 0.12, 7.5, fill=SAFFRON)
rect(sl, 0.12, 0.0, 0.06, 7.5, fill=MAROON)

# Mandala corners — bottom only so they don't touch the text block (y≥6.15)
corner_mandala(sl, 0.42, 7.25, max_r=0.62, rings=3)  # bottom-left, clear of text
corner_mandala(sl, 12.5, 6.65, max_r=1.3, rings=5)   # bottom-right

logo_bar(sl)
ai_gem(sl)

txt(sl, "Proposal by", 1.0, 1.5, 6.0, 0.6, size=22, italic=True, color=ROSE_GOLD)
txt(sl, "SAHELI", 0.9, 2.1, 9.0, 1.9, size=104, bold=True, color=CREAM)
txt(sl, "सहेली", 1.0, 3.9, 5.0, 0.85, size=32, bold=True, color=SAFFRON)

gold_line(sl, 1.0, 4.85, 7.5, thick=0.04)

txt(sl,
    "India's First Voice-Native, Stealth Evidence Vault\n"
    "& Rights Companion for Women on Shared Phones",
    1.0, 5.0, 8.5, 1.1, size=16, italic=True, color=CREAM)
txt(sl,
    "WitchHunt 2026  |  HopeWorks × AI4India × Skills Café\n"
    "Track: Health & Well-Being   ✦   Cross-track: Smart Cities, Education",
    1.0, 6.15, 9.0, 0.8, size=11, color=ROSE_GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Crisis in Numbers
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Title band
rect(sl, 0.0, 0.0, 13.33, 1.0, fill=MAROON)
# ✦ No random gold blobs — replaced with small diamond corner marks
corner_mandala(sl, 12.9, 0.5, max_r=0.7, rings=3)
corner_mandala(sl, 0.43, 0.5, max_r=0.7, rings=3)
logo_bar(sl)
ai_gem(sl)
txt(sl, "CRISIS IN NUMBERS", 1.0, 0.18, 11.3, 0.65,
    size=32, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

# Stat cards — taller, no gold blobs, accent colour top-bar per card
stats = [
    ("43%",  "of Indian women\ndo not own a smartphone", SAFFRON),
    ("58%",  "of rural women\nuse shared devices",        LOTUS),
    ("67%",  "harassment cases\ngo unreported",           TEAL),
    ("26%",  "DV conviction rate\ndue to evidence gaps",  ROSE_GOLD),
    ("74%",  "Nirbhaya Fund\nstill sits unspent",         GOLD_LIGHT),
]
xs = [0.2, 2.82, 5.44, 8.06, 10.68]
card_sub = [
    "Existing apps require\na personal device",
    "Shared phones mean\nconstant surveillance",
    "Fear + no private channel\n= forced silence",
    "Without evidence,\ncases collapse",
    "Systemic failure\nat the last mile",
]
for (num, label, accent), lx, sub in zip(stats, xs, card_sub):
    rect(sl, lx, 1.1, 2.42, 4.55, fill=RGBColor(0x12,0x06,0x28), line=accent, lw=Pt(1.2))
    rect(sl, lx, 1.1, 2.42, 0.22, fill=accent)      # solid colour top bar
    txt(sl, num, lx, 1.52, 2.42, 1.0,
        size=54, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(sl, label, lx+0.12, 2.65, 2.2, 1.1, size=12,
        color=CREAM, align=PP_ALIGN.CENTER, wrap=True)
    # thin divider
    rect(sl, lx+0.3, 3.85, 1.82, 0.02, fill=RGBColor(0x40,0x30,0x55))
    # context sub-text fills lower third
    txt(sl, sub, lx+0.12, 3.95, 2.2, 0.95, size=10.5,
        color=RGBColor(0xAA,0xA0,0xCC), align=PP_ALIGN.CENTER, wrap=True, italic=True)
    # small accent ring at bottom of card
    oval(sl, lx+1.0, 5.2, 0.4, 0.4, line=accent, lw=Pt(1))

gold_line(sl, 0.2, 5.78, 12.93, thick=0.035)
rect(sl, 0.2, 5.83, 12.93, 1.35, fill=DEEP_CARD, line=LOTUS, lw=Pt(1))
txt(sl,
    "Every existing solution assumes private smartphones, stable internet, and high digital literacy.\n"
    "180M+ Indian women have none of these.  Saheli is built for that gap.",
    0.45, 5.93, 12.43, 0.95, size=14, italic=True,
    color=CREAM, align=PP_ALIGN.CENTER)
txt(sl,
    "Sources: NSSO  ✦  IAMAI  ✦  GSMA Connected Women 2024  ✦  NCRB  ✦  NARI  ✦  NALSA",
    0.5, 7.06, 12.3, 0.28, size=9, color=ROSE_GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Left panel (full height)
rect(sl, 0.0, 0.0, 3.5, 7.5, fill=MAROON)
rich_mandala(sl, 1.75, 2.5, max_r=2.6, rings=5)   # fills top half of panel
# second accent ring in lower half so panel isn't empty below mandala
petal_ring(sl, 1.75, 5.5, r=0.9, n=10, color=GOLD_LIGHT)
corner_mandala(sl, 1.75, 5.5, max_r=1.0, rings=3)

txt(sl, "PROBLEM\nSTATEMENT", 0.18, 5.6, 3.2, 1.4, size=24, bold=True, color=CREAM)
gold_line(sl, 0.2, 5.55, 3.0, thick=0.04)
txt(sl, "The Invisible Evidence Gap", 0.18, 7.05, 3.1, 0.38,
    size=10, italic=True, color=GOLD_LIGHT)

logo_bar(sl)
ai_gem(sl)

# Vertical timeline (clear of left panel)
rect(sl, 4.18, 0.18, 0.05, 7.1, fill=RGBColor(0x40,0x30,0x60))
for y in [0.9, 2.5, 4.0, 5.55]:
    oval(sl, 4.07, y-0.12, 0.25, 0.25, fill=SAFFRON)

sections = [
    ("Problem Statement",          GOLD_LIGHT,
     "Shared phones make safety and rights invisible for 180M+ Indian women."),
    ("Problem Statement Expansion", SAFFRON,
     "Most women's safety tools assume individual smartphone ownership, stable internet, and English "
     "literacy. 43% of Indian women don't own a phone. 58% in rural areas share devices with those "
     "they need protection from. Any visible app is a danger."),
    ("How it manifests for the user", LOTUS,
     "Safety apps are visible on shared devices — a DV survivor cannot install them. Complex UI excludes "
     "WhatsApp-first users. SOS-only tools miss legal aid, evidence documentation, and scheme access."),
    ("How it impacts the user or client", ROSE_GOLD,
     "Abuse stays undocumented and cases collapse. Women miss legal aid, government schemes, and care. "
     "Institutions exist — OSCs, NALSA, ASHA — but the last-mile connection fails the women who need it most."),
]
ty = 0.38
for title, tcol, body in sections:
    txt(sl, title, 4.5, ty, 8.65, 0.36, size=12.5, bold=True, color=tcol)
    txt(sl, body,  4.5, ty+0.36, 8.65, 1.1, size=11.5, color=CREAM, wrap=True)
    ty += 1.6

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Project Details
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Right panel — circular Warli tribal art (mirrors mandala geometry, real Indian art)
img(sl, "warli_painting.jpg", 6.85, 0.0, 6.48, 7.5)
# Subtle teal petal overlay on top of photo
petal_ring(sl, 10.09, 3.75, r=2.9, n=20, color=TEAL)

logo_bar(sl)
ai_gem(sl)

# Card (left half — clear of badge which ends at x≈0.8)
rect(sl, 0.28, 0.45, 6.38, 6.9,
     fill=DEEP_CARD, line=SAFFRON, lw=Pt(1.5))
gold_line(sl, 0.28, 0.45, 6.38, thick=0.05)

txt(sl, "Project Details", 0.55, 0.58, 5.8, 0.65, size=30, bold=True, color=CREAM)
gold_line(sl, 0.55, 1.22, 5.8, thick=0.03)
txt(sl, "SAHELI  (सहेली)", 0.55, 1.35, 5.8, 0.48, size=20, bold=True, color=SAFFRON)

txt(sl,
    "A voice-native, stealth-first evidence vault and rights companion "
    "disguised as a devotional app. One sentence in Hindi triggers evidence "
    "capture, legal guidance, SOS alerts, scheme eligibility, and health "
    "screening — all offline-capable on a ₹1,000 phone.",
    0.55, 1.9, 5.9, 1.65, size=12.5, color=CREAM, wrap=True)

gold_pill(sl, "DEMOGRAPHIC IT IMPACTS", 0.55, 3.68, 3.5)
txt(sl,
    "Women aged 18–45 in tier-2/3 India: shared-phone households, "
    "rural users with feature phones, low-literacy survivors navigating "
    "abuse, harassment, and health gaps.",
    0.55, 4.15, 5.9, 1.15, size=12.5, color=CREAM, wrap=True)

gold_pill(sl, "COMPETITION TRACK", 0.55, 5.42, 2.8)
txt(sl, "Health & Well-Being  ✦  Smart Cities  ✦  Education",
    0.55, 5.88, 5.9, 0.38, size=11.5, color=ROSE_GOLD)

gold_line(sl, 0.28, 7.3, 6.38, thick=0.04)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Saheli Design Principles
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Right panel — peacock lotus rangoli (grace + design principles)
img(sl, "sankranthi.jpg", 6.85, 0.0, 6.48, 7.5)
petal_ring(sl, 10.09, 3.75, r=2.9, n=20, color=LOTUS)

logo_bar(sl)
ai_gem(sl)

# FIX: title starts at x=1.0 to clear the AI badge (badge ends at x≈0.8)
txt(sl, "Saheli Design Principles",
    1.0, 0.22, 5.6, 0.62, size=26, bold=True, color=CREAM)
gold_line(sl, 1.0, 0.86, 5.6, thick=0.04)
txt(sl, "Built for real humans, not hypothetical users.",
    1.0, 0.96, 5.6, 0.36, size=12, italic=True, color=SAFFRON)

rect(sl, 0.28, 1.42, 6.38, 6.0, fill=DEEP_CARD, line=LOTUS, lw=Pt(1))

principles = [
    "Voice-first: one conversation replaces every form and every menu",
    "Stealth-first: invisible on shared phones, hidden inside a devotional app",
    "Offline-first: works on ₹1,000 phones with zero data plan",
    "Build for need and necessity — not comfort and convenience",
    "Assume constraints, not perfect conditions",
    "Design for simplicity, empathy, and dignity",
    "Aim for solutions that actually change someone's day-to-day reality",
    "Reach the woman with the shared phone first — she is the real user",
]
ty = 1.57
for p in principles:
    txt(sl, "✦", 0.42, ty, 0.38, 0.52, size=12, bold=True, color=SAFFRON)
    txt(sl, p,   0.82, ty, 5.7, 0.52, size=12.5, color=CREAM, wrap=True)
    ty += 0.68

# ══════════════════════════════════════════════════════════════════════════════
# PERSONA SLIDE BUILDER — panels fill to y≈6.75, journey strip at bottom
# ══════════════════════════════════════════════════════════════════════════════
def persona_slide(prs, num, name, age_loc, desc, traits, triggers, problems,
                  goals, journey, photo=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sl, 0x0B, 0x07, 0x24)
    logo_bar(sl)

    SIDEBAR_TOP = 0.78
    JOURNEY_H   = 0.8       # height of journey strip at bottom
    GRID_BOTTOM = 7.5 - JOURNEY_H - 0.08  # ≈ 6.62

    # ── Title band ────────────────────────────────────────────────────────────
    rect(sl, 0.0, 0.0, 13.33, SIDEBAR_TOP, fill=MAROON)
    corner_mandala(sl, 12.9, 0.38, max_r=0.65, rings=3)
    corner_mandala(sl, 0.43, 0.38, max_r=0.65, rings=3)
    txt(sl, f"USER PERSONA {num}", 0.5, 0.1, 12.3, 0.6,
        size=28, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

    # ── Left sidebar ──────────────────────────────────────────────────────────
    rect(sl, 0.0, SIDEBAR_TOP, 2.55, 7.5 - SIDEBAR_TOP,
         fill=RGBColor(0x10, 0x05, 0x28))
    if photo:
        img(sl, photo, 0.2, SIDEBAR_TOP + 0.1, 2.15, 2.85)
        rect(sl, 0.2, SIDEBAR_TOP + 0.1, 2.15, 2.85, line=SAFFRON, lw=Pt(1.5))
    else:
        rect(sl, 0.2, SIDEBAR_TOP + 0.1, 2.15, 2.85,
             fill=RGBColor(0x28, 0x18, 0x48), line=SAFFRON, lw=Pt(1))
        txt(sl, "[ Photo ]", 0.2, SIDEBAR_TOP + 1.0, 2.15, 0.4,
            size=9, color=ROSE_GOLD, align=PP_ALIGN.CENTER, italic=True)
    name_pill_t = SIDEBAR_TOP + 3.05
    rect(sl, 0.2, name_pill_t, 2.15, 0.52, fill=LOTUS)
    txt(sl, name.upper(), 0.2, name_pill_t + 0.05, 2.15, 0.42,
        size=17, bold=True, color=MIDNIGHT, align=PP_ALIGN.CENTER)
    txt(sl, age_loc, 0.22, name_pill_t + 0.62, 2.1, 1.8,
        size=11, color=CREAM, wrap=True)
    corner_mandala(sl, 1.27, GRID_BOTTOM - 0.4, max_r=0.7, rings=3)

    # ── Right content grid ────────────────────────────────────────────────────
    GRID_LEFT = 2.6
    GRID_TOP  = SIDEBAR_TOP + 0.04
    ROW1_H    = 2.35
    ROW2_H    = GRID_BOTTOM - GRID_TOP - ROW1_H - 0.08   # ≈ 3.35"

    label_cols = [SAFFRON, GOLD_LIGHT, LOTUS, SAFFRON, TEAL]
    bgs        = [RGBColor(0x1A,0x06,0x38), RGBColor(0x12,0x08,0x28),
                  RGBColor(0x1A,0x06,0x38), RGBColor(0x12,0x08,0x28),
                  RGBColor(0x1A,0x06,0x38)]

    def panel(label, content, pl, pt, pw, ph, lc, bg):
        rect(sl, pl, pt, pw, ph, fill=bg, line=lc, lw=Pt(1))
        rect(sl, pl, pt, pw, 0.38, fill=lc)
        txt(sl, label, pl+0.1, pt+0.05, pw-0.15, 0.3,
            size=9, bold=True, color=MIDNIGHT)
        txt(sl, content, pl+0.14, pt+0.45, pw-0.25, ph-0.52,
            size=11, color=CREAM, wrap=True)

    # Row 1: About (wider) + Traits
    about_w  = 5.85
    traits_w = 13.33 - GRID_LEFT - about_w - 0.08
    panel("ABOUT THE USER", desc,
          GRID_LEFT, GRID_TOP, about_w, ROW1_H, label_cols[0], bgs[0])
    panel("TRAITS", "\n".join(f"✦ {t}" for t in traits),
          GRID_LEFT + about_w + 0.08, GRID_TOP, traits_w, ROW1_H,
          label_cols[1], bgs[1])

    # Row 2: Triggers / Problems / Goals — fills remaining space
    ROW2_TOP = GRID_TOP + ROW1_H + 0.08
    avail_w  = 13.33 - GRID_LEFT - 0.16
    col_w    = avail_w / 3 - 0.05
    for i, (label, items) in enumerate([
        ("TRIGGER MOMENTS", triggers),
        ("PROBLEMS",        problems),
        ("GOALS AND NEEDS", goals),
    ]):
        px = GRID_LEFT + i * (col_w + 0.09)
        panel(label, "\n".join(f"✦ {item}" for item in items),
              px, ROW2_TOP, col_w, ROW2_H, label_cols[i+2], bgs[i+2])

    # ── Journey strip (fills bottom of slide) ────────────────────────────────
    JOURNEY_TOP = GRID_BOTTOM + 0.06
    rect(sl, 0.0, JOURNEY_TOP, 13.33, JOURNEY_H + 0.04,
         fill=RGBColor(0x1A, 0x06, 0x2E), line=SAFFRON, lw=Pt(0.8))
    txt(sl, "✦  JOURNEY SNAPSHOT", 0.35, JOURNEY_TOP + 0.06, 2.5, 0.28,
        size=8, bold=True, color=SAFFRON)
    txt(sl, journey, 2.85, JOURNEY_TOP + 0.05, 10.3, 0.65,
        size=10.5, color=CREAM, italic=True, wrap=True)

    return sl

# ── Persona 1: Priya ──────────────────────────────────────────────────────────
persona_slide(
    prs, 1, "Priya",
    "• 35 years old\n• Lucknow\n• Shared-phone\n  mother\n• 2 children",
    "Priya lives with an abusive husband who controls the phone, the money, "
    "and what apps stay installed. She needs to document abuse, understand "
    "her rights, and find help — without ever appearing to do so.",
    ["Cautious", "Resourceful", "Protective", "Hindi-first", "Privacy-sensitive"],
    ["Phone use is monitored", "Needs proof without a visible app",
     "Wants rights explained safely"],
    ["No secure evidence storage", "Limited DV Act knowledge",
     "Cannot openly ask for help"],
    ["Capture proof invisibly", "Understand maintenance, custody & property",
     "Reach OSC or trusted contacts"],
    "Week 1: friend installs 'Daily Mantras', husband approves  →  Week 4: 7 evidence items SHA-256 hashed  →  Week 6: stealth SOS configured  →  Week 8: OSC contacted with full case file",
    photo="priya.jpg",
)

# ── Persona 2: Sunita ─────────────────────────────────────────────────────────
persona_slide(
    prs, 2, "Sunita",
    "• 29 years old\n• Nashik\n• Teacher\n• Basic smartphone",
    "Sunita has lived with severe period pain for 13 years and calls it "
    "normal because care is expensive and no one explained otherwise. "
    "She needs structured screening, affordable medicine, and a Hindi "
    "doctor-ready summary.",
    ["Resilient", "Underdiagnosed", "Budget-conscious", "Working daily", "Seeks relief"],
    ["Pain is normalized at home", "Doctor visits feel expensive",
     "Needs simple Hindi guidance"],
    ["No symptom history to share", "No structured screening flow",
     "Affordable care hard to find"],
    ["Track symptoms over time", "Access Jan Aushadhi (₹3/day vs ₹30)",
     "Carry a structured Doctor Card"],
    "ASHA shares link → 'Pet mein dard' → PCOS screening positive → exercise protocol → Doctor Card generated → Jan Aushadhi ₹3/day vs ₹30 → first doctor visit in 13 years",
    photo="sunita.jpg",
)

# ── Persona 3: Meena ──────────────────────────────────────────────────────────
persona_slide(
    prs, 3, "Meena",
    "• 42 years old\n• Varanasi\n• Lava feature phone\n• No internet",
    "Meena has no app-install capacity, no data plan, and no habit of "
    "digital healthcare. If support cannot reach her through SMS it does "
    "not reach her at all. She has never been screened for any health "
    "condition in her life.",
    ["Low-literacy", "Rural", "Persistent", "Never screened", "Trusts ASHA workers"],
    ["No data plan or app store", "Needs SMS-first guidance",
     "Lives far from formal care"],
    ["No installable app path", "Screening info never reaches her",
     "High-risk signs go unchecked"],
    ["Screen through SMS/USSD", "Get nearest PHC guidance offline",
     "Receive follow-up via ASHA"],
    "SMS 'SAHELI' to shortcode → 5-SMS cervical risk screen → high risk flagged → 'Free test, PHC 2 km' → ASHA follow-up scheduled → first health screening in 42 years",
    photo="meena.jpg",
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Right decoration panel (narrower — more room for Q&A)
rect(sl, 8.6, 0.0, 4.73, 7.5, fill=RGBColor(0x10, 0x05, 0x28))
rich_mandala(sl, 10.97, 3.75, max_r=3.5, rings=6)
petal_ring(sl, 10.97, 3.75, r=2.2, n=14, color=TEAL)

logo_bar(sl)
ai_gem(sl)

# FIX: title at x=1.0, clear of badge
txt(sl, "Our Solution", 1.0, 0.22, 7.0, 0.75, size=38, bold=True, color=CREAM)
gold_line(sl, 1.0, 0.98, 7.4, thick=0.04)
txt(sl,
    "A stealth, voice-native rights and evidence companion\nfor women on shared phones.",
    1.0, 1.1, 7.0, 0.85, size=14, bold=True, color=SAFFRON)

# Left graphic panel — real photo: colorful Indian women walking (community)
img(sl, "r_varanasi_10.jpg", 0.2, 2.1, 3.45, 5.25)
# Thin lotus-pink border frame around the photo
rect(sl, 0.2, 2.1, 3.45, 5.25, line=LOTUS, lw=Pt(1.5))

# Vertical timeline bar
rect(sl, 3.85, 2.1, 0.05, 5.2, fill=RGBColor(0x40,0x30,0x60))
for y in [2.4, 3.25, 4.1, 4.95, 5.8, 6.65]:
    oval(sl, 3.74, y-0.12, 0.24, 0.24, fill=SAFFRON)

questions = [
    ("What do you plan to build?",
     "Voice-first stealth app: evidence vault + rights guide + SOS + scheme finder + health journal."),
    ("Who is it for?",
     "180M+ women aged 18–45 in tier-2/3 India — shared phones, feature phones, offline users."),
    ("What problem does it address?",
     "The invisible evidence gap: women can't document abuse or access rights on a monitored device."),
    ("How would it work in the real world?",
     "One sentence in Hindi → Saheli captures evidence, explains rights, or triggers SOS offline."),
    ("What impact would it create?",
     "Documented cases, legal empowerment, reduced burnout — at scale via ASHA and SHG networks."),
    ("Why is it worth building?",
     "300M+ women are excluded by every existing platform. Saheli closes that gap."),
]
label_colors = [GOLD_LIGHT, LOTUS, TEAL, SAFFRON, ROSE_GOLD, LOTUS]

ty = 2.18
for (q, a), lc in zip(questions, label_colors):
    txt(sl, q, 4.15, ty, 4.2, 0.3, size=12, bold=True, color=lc)
    txt(sl, a, 4.15, ty+0.3, 4.2, 0.58, size=10.5, color=CREAM, wrap=True)
    ty += 0.85

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Our Solution Continue (Five Capabilities)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

rect(sl, 0.0, 0.0, 13.33, 0.98, fill=RGBColor(0x18,0x06,0x36))
gold_line(sl, 0.0, 0.98, 13.33, thick=0.04)
logo_bar(sl)
ai_gem(sl)

# FIX: title at x=1.0
txt(sl, "Our Solution Continue", 1.0, 0.1, 9.0, 0.62, size=26, bold=True, color=CREAM)
txt(sl, "Five capabilities. One conversation.", 1.0, 0.66, 9.0, 0.28,
    size=13, bold=True, color=SAFFRON)

# Graphic placeholder with rich mandala
rect(sl, 0.2, 1.08, 2.65, 3.8, fill=RGBColor(0x0F,0x05,0x28), line=LOTUS, lw=Pt(1))
rich_mandala(sl, 1.52, 2.98, max_r=1.9, rings=4)

caps = [
    (SAFFRON,   "01", "Evidence Vault  (सबूत तिजोरी)",
     "SHA-256 hashed photos, recordings, documents, GPS timestamps. "
     "Aligned to DV Act 2005, IPC 498A, POSH Act. Cover app shows nothing."),
    (LOTUS,     "02", "Know Your Rights  (अधिकार)",
     "RAG on PWDVA, POSH, Hindu Succession, maintenance, property. "
     "Voice delivery in Hindi/regional via Bulbul V3 TTS. Divorce calculator."),
    (TEAL,      "03", "Safety Net  (सुरक्षा)",
     "Stealth SOS from mala counter. GPS to 3 contacts, offline SMS fallback. "
     "733+ One Stop Centre locator. Cached helplines: 181, 112, 1091."),
    (GOLD_LIGHT,"04", "Entitlements  (योजना)",
     "20+ schemes: PMMVY, Ayushman Bharat, Mudra, Sukanya Samriddhi. "
     "OCR scan of Aadhaar/Ration Card → auto eligibility via Sarvam Vision."),
    (ROSE_GOLD, "05", "Health Journal  (सेहत)",
     "Voice EPDS, PCOS, cervical risk screening. Evidence-based stretching "
     "(Frontiers Med 2025). Doctor Card generator. ABHA integration."),
]
ty = 1.08
for num_c, num, title, body in caps:
    rect(sl, 3.05, ty, 10.05, 1.25, fill=DEEP_CARD, line=num_c, lw=Pt(1.2))
    rect(sl, 3.05, ty, 0.52, 1.25, fill=num_c)       # solid left-edge colour bar
    # wrap=False + wider box prevents "01" splitting to "0\n1"
    txt(sl, num, 3.06, ty+0.38, 0.5, 0.42, size=16, bold=True, color=MIDNIGHT,
        align=PP_ALIGN.CENTER, wrap=False)
    txt(sl, title, 3.65, ty+0.1, 9.2, 0.36, size=13, bold=True, color=CREAM)
    txt(sl, body,  3.65, ty+0.5, 9.2, 0.68, size=10.5, color=ROSE_GOLD, wrap=True)
    ty += 1.28

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Use of Artificial Intelligence
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x08, 0x03, 0x1A)

# Left panel — real lotus photo (symbol of transformation / Saheli's emblem)
img(sl, "r_lotus_pink.jpg", 0.0, 0.0, 4.8, 7.5)
# Mandala ring overlay on top of photo for mystical layer
petal_ring(sl, 2.4, 3.75, r=2.2, n=12, color=LOTUS)
logo_bar(sl)
ai_gem(sl)

# Small "USE OF" label anchored above the title on the left of content area
txt(sl, "USE OF", 4.8, 0.32, 2.0, 0.3,
    size=10, bold=True, color=ROSE_GOLD, align=PP_ALIGN.LEFT)

# Headline — moved down slightly to sit cleanly below the tag
txt(sl, "Artificial\nIntelligence", 4.95, 0.94, 8.2, 2.1, size=50, bold=True, color=CREAM)
gold_line(sl, 4.95, 3.08, 8.1, thick=0.04)

txt(sl,
    "AI is the interface layer, not the product. It converts complex legal, health, "
    "and scheme workflows into one sentence in Hindi. The outcomes — evidence, rights, "
    "safety — are real whether or not AI drafted the response.",
    4.95, 3.22, 8.1, 1.3, size=13, italic=True, color=CREAM, wrap=True)

# FIX: tech stack starts immediately after body (no 1.5" gap)
ai_stack = [
    ("Saaras V3 STT",  "Hindi-first voice capture across 22+ Indian languages"),
    ("Bulbul V3 TTS",  "Empathetic spoken guidance — not robotic readout"),
    ("Sarvam LLM",     "Domain-specific RAG on Bare Acts, WHO protocols, NALSA"),
    ("Sarvam Vision",  "OCR on Aadhaar/Ration Card for scheme eligibility"),
    ("Edge AI",        "Full offline inference on low-RAM Android devices"),
]
ty = 4.6
for label, desc_text in ai_stack:
    txt(sl, f"✦  {label}", 4.95, ty, 2.55, 0.38, size=11, bold=True, color=SAFFRON)
    txt(sl, f"—  {desc_text}", 7.45, ty, 5.55, 0.38, size=11, color=CREAM)
    ty += 0.46

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Challenges and Risks
# FIX: top strip 2.9 → 1.55"; challenges no longer cut off at bottom
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x08, 0x03, 0x1A)

# Reduced top strip (1.55" instead of 2.9")
rect(sl, 0.0, 0.0, 13.33, 1.55, fill=RGBColor(0x14,0x08,0x30))
petal_ring(sl, 3.5, 0.77, r=0.65, n=10, color=SAFFRON)
petal_ring(sl, 9.8, 0.77, r=0.65, n=10, color=LOTUS)
corner_mandala(sl, 0.55, 0.77, max_r=0.65, rings=3)
corner_mandala(sl, 12.75, 0.77, max_r=0.65, rings=3)
txt(sl, "Challenges  ✦  Risks  ✦  Mitigations", 4.5, 0.6, 4.3, 0.4,
    size=10, color=ROSE_GOLD, align=PP_ALIGN.CENTER, italic=True)

logo_bar(sl)
gold_line(sl, 0.0, 1.55, 13.33, thick=0.04)

# Title + left-zone content (fills the dead x=0→5 zone)
txt(sl, "Challenges\nand Risks", 0.4, 1.72, 4.6, 1.55, size=36, bold=True, color=CREAM)
gold_line(sl, 0.4, 3.3, 4.4, thick=0.03)

# Fill left area below title with a key principle card
rect(sl, 0.3, 3.42, 4.6, 3.55, fill=RGBColor(0x14,0x06,0x2C), line=MAROON, lw=Pt(1))
txt(sl, "Key Design Principle", 0.55, 3.56, 4.1, 0.32, size=10, bold=True, color=SAFFRON)
txt(sl,
    "\"Build as if every\nfeature could get\nsomeone hurt.\"",
    0.55, 3.95, 4.0, 1.1, size=13, italic=True, color=CREAM, wrap=True)
# Thin divider inside card
rect(sl, 0.55, 5.12, 3.8, 0.025, fill=MAROON)
txt(sl,
    "Safety is not a feature —\nit is the constraint that\nshapes every decision.",
    0.55, 5.2, 4.0, 0.88, size=12, italic=True, color=ROSE_GOLD, wrap=True)
# Decorative petal ring at bottom of card
petal_ring(sl, 2.6, 6.6, r=0.28, n=8, color=MAROON)
corner_mandala(sl, 2.6, 6.6, max_r=0.4, rings=2)
txt(sl, "✦  सुरक्षा  ✦", 0.55, 6.75, 4.0, 0.28,
    size=10, color=MAROON, align=PP_ALIGN.CENTER, italic=True)

# 4 challenges — with reduced spacing so all fit
challenges = [
    (SAFFRON,   "Trust & Safety",
     "AI must never sound like legal or medical advice. Guardrails and escalation paths required at every response."),
    (LOTUS,     "Stealth Failure",
     "Any visible trace on a shared phone can endanger the user. Quick exit, hidden notifications, and cover-app integrity are non-negotiable."),
    (TEAL,      "Localization Drift",
     "Dialect accuracy and cached legal/health knowledge can become stale. Curated update cycles are built into the roadmap."),
    (ROSE_GOLD, "Last-Mile Reach",
     "SMS/USSD path must be maintained so feature-phone users are never cut off from safety-critical features."),
]
ty = 1.65
for lc, title, body in challenges:
    txt(sl, f"✦  {title}:", 5.2, ty, 7.8, 0.36, size=13, bold=True, color=lc)
    txt(sl, body,           5.2, ty+0.36, 7.8, 0.78, size=12, color=CREAM, wrap=True)
    ty += 1.18    # reduced from 1.1 — but now within 1.55+4*1.18=6.27" total, well within 7.5"

# Footer bar
rect(sl, 0.0, 7.1, 13.33, 0.38, fill=MAROON)
txt(sl,
    "The product only works if privacy, accuracy, and last-mile delivery outperform every "
    "existing safety app — simultaneously.",
    0.35, 7.13, 12.63, 0.3, size=10, italic=True, color=CREAM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Technical Request for Facilitated Datasets
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

rect(sl, 0.0, 0.0, 13.33, 1.0, fill=RGBColor(0x18,0x06,0x36))
gold_line(sl, 0.0, 1.0, 13.33, thick=0.04)
logo_bar(sl)
ai_gem(sl)

# FIX: title starts at x=0.9 to clear AI badge (badge ends x≈0.8)
txt(sl, "Technical Request for Facilitated Datasets",
    0.9, 0.12, 12.0, 0.62, size=22, bold=True, color=CREAM)
txt(sl, "Categorized parameters for AI4India facilitated data access:",
    0.9, 0.7, 12.0, 0.3, size=11, italic=True, color=SAFFRON)

# Symmetric column divider at x=6.66 (half of 13.33)
rect(sl, 6.62, 1.08, 0.05, 6.2, fill=RGBColor(0x40,0x30,0x60))
COL_L = 0.4
COL_R = 6.82
COL_W = 6.05  # symmetric on both sides

left_secs = [
    ("Specification of Required Data", SAFFRON, [
        ("Data Type:",
         "Indic speech from women in low-resource dialects; verified OSC/PHC/ASHA directories; anonymized women's health records."),
        ("Volume & Depth:",
         "10,000+ annotated audio clips; 5-year scheme data; district-level service coverage."),
        ("Granularity:",
         "Village-level service locations; daily scheme update feeds."),
    ]),
    ('The "Ask" vs. "Leveraged" Mapping', LOTUS, [
        ("What we have:",
         "Bare Acts, NALSA, WHO/UNICEF protocols, open IndiaAI datasets, Bhashini corpora."),
        ("What we need:",
         "Facilitated dialect voice data and verified service directories — the missing piece for rural accuracy."),
    ]),
]
right_secs = [
    ("Technical Justification", TEAL, [
        ("Accuracy:",
         "Improves STT for rural voices and code-mixed Hindi queries."),
        ("Validation:",
         "Grounds EPDS, PCOS, and cervical screening in verified clinical data."),
        ("Edge-Case Handling:",
         "Dialect data ensures Meena-tier users are never excluded."),
    ]),
    ("Environment & Security Requirements", ROSE_GOLD, [
        ("Access Method:",
         "Secure API, sandbox, or encrypted batch download."),
        ("Security:",
         "Minimum-retention storage; DPDP-compliant pipeline; no PII retained in prompts."),
    ]),
]

def render_two_col_sec(secs, col_x):
    ty = 1.12
    for sec_title, sc, items in secs:
        txt(sl, sec_title, col_x, ty, COL_W, 0.36, size=12, bold=True, color=sc)
        gold_line(sl, col_x, ty+0.36, COL_W, thick=0.022)
        ty += 0.47
        for label, body in items:
            txt(sl, label, col_x,       ty, 1.45, 0.28, size=10, bold=True, color=GOLD_LIGHT)
            txt(sl, body,  col_x+1.48,  ty, COL_W-1.52, 0.82, size=10, color=CREAM, wrap=True)
            ty += 0.76
        ty += 0.18

render_two_col_sec(left_secs,  COL_L)
render_two_col_sec(right_secs, COL_R)

txt(sl, "Request focus: voice accuracy  ✦  service routing  ✦  health validation",
    0.4, 7.12, 12.5, 0.3, size=10, italic=True, color=ROSE_GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Anticipated Results
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

rect(sl, 0.0, 0.0, 13.33, 0.98, fill=RGBColor(0x18,0x06,0x36))
gold_line(sl, 0.0, 0.98, 13.33, thick=0.04)
logo_bar(sl)
ai_gem(sl)

# FIX: title centred but clear of badge (badge is left, title is centred at x≥1)
txt(sl, "Anticipated Results", 1.0, 0.1, 11.33, 0.62,
    size=30, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
txt(sl, "Year 1 targets anchored to documented action, not vanity metrics.",
    1.0, 0.7, 11.33, 0.28, size=12, italic=True, color=SAFFRON, align=PP_ALIGN.CENTER)

# Graphic — real rangoli festival art (light+celebration = impact)
img(sl, "r_rangoli_festival.jpg", 0.2, 1.08, 2.9, 4.0)
rect(sl, 0.2, 1.08, 2.9, 4.0, line=TEAL, lw=Pt(1.5))

# Timeline bar
rect(sl, 3.45, 1.08, 0.05, 6.1, fill=RGBColor(0x40,0x30,0x60))
txt(sl, "Saheli Impact Model", 3.65, 1.1, 9.5, 0.42, size=15, bold=True, color=CREAM)

results = [
    (SAFFRON,   "Primary beneficiary  (Priya — DV Survivor)",
     "Estimated 40% increase in evidence-backed DV filings. OSC case completion rate improved through guided documentation."),
    (LOTUS,     "Secondary beneficiary  (Sunita — Health Gap)",
     "First structured screening for 50,000+ women in Year 1. Jan Aushadhi saves ~₹27/day per user."),
    (TEAL,      "Feature-phone reach  (Meena — Rural / Offline)",
     "SMS/USSD path onboards users with zero internet. ASHA network drives last-mile coverage at no additional infrastructure cost."),
    (ROSE_GOLD, "Solution scalability — operational and technical",
     "ASHA/SHG distribution requires zero new infrastructure. Edge inference scales to ₹1,000 devices. B2G + B2B revenue sustains the free core."),
]
ty = 1.65
for lc, title, body in results:
    oval(sl, 3.35, ty+0.04, 0.22, 0.22, fill=lc)
    txt(sl, title, 3.72, ty, 9.3, 0.32, size=13, bold=True, color=lc)
    txt(sl, body,  3.72, ty+0.33, 9.3, 0.68, size=11, color=CREAM, wrap=True)
    ty += 1.05

# Distribution & sustainability strip (fills bottom empty space)
DIST_TOP = ty + 0.05
if DIST_TOP < 6.2:   # safety guard
    DIST_TOP = 6.2
rect(sl, 3.3, DIST_TOP, 9.85, 1.12,
     fill=RGBColor(0x14,0x06,0x2C), line=TEAL, lw=Pt(0.8))
txt(sl, "Distribution & Sustainability",
    3.5, DIST_TOP + 0.08, 9.4, 0.3, size=10, bold=True, color=TEAL)
dist_items = [
    "ASHA network (1M+) — primary distribution at zero infrastructure cost",
    "SHG network (9M) — group sessions, leader runs Saheli for 10–15 women",
    "Revenue: B2G (ASHA/OSC tools)  ✦  B2B (POSH compliance)  ✦  Grants — core always free",
]
txt(sl, "   ✦   ".join(dist_items),
    3.5, DIST_TOP + 0.42, 9.55, 0.62, size=10, color=CREAM, italic=True, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(sl, 0x0B, 0x07, 0x24)

# Background rings (subtle vignette on left)
oval(sl, -1.2, -1.2, 8.0, 9.8, line=RGBColor(0x20,0x10,0x45), lw=Pt(1.5))
oval(sl, -0.3, -0.3, 6.0, 8.0, line=RGBColor(0x30,0x15,0x55), lw=Pt(1))

# Right panel — Warli painting bookends with the cover slide
img(sl, "r_warli_jivya.jpg", 7.5, 0.0, 5.83, 7.5)
# Colour-ring overlay to blend into palette
petal_ring(sl, 11.2, 3.75, r=2.55, n=18, color=SAFFRON)
petal_ring(sl, 11.2, 3.75, r=1.55, n=12, color=LOTUS)

# Left accent stripe
rect(sl, 0.0, 0.0, 0.12, 7.5, fill=SAFFRON)
rect(sl, 0.12, 0.0, 0.06, 7.5, fill=MAROON)

logo_bar(sl)
ai_gem(sl)

txt(sl, "THANK YOU", 0.7, 1.55, 7.8, 1.85, size=80, bold=True, color=CREAM)
gold_line(sl, 0.7, 3.45, 7.2, thick=0.05)
txt(sl, "for your time and attention", 0.7, 3.62, 7.2, 0.7, size=22, color=CREAM)
txt(sl, "Saheli  ✦  सहेली", 0.7, 4.52, 7.2, 0.55, size=18, bold=True, color=SAFFRON)
txt(sl,
    "one conversation for every right, every scheme, and every tool",
    0.7, 5.15, 7.2, 0.55, size=13, italic=True, color=ROSE_GOLD)
txt(sl,
    "WitchHunt 2026  |  HopeWorks Foundation × AI4India × Skills Café\n"
    "Track: Health & Well-Being",
    0.7, 6.3, 7.2, 0.8, size=10, color=RGBColor(0x80,0x70,0x90))

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/garudev/dev/saheli/Saheli_Solution_Proposal.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
