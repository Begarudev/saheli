from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Solution Proposal - Template.pptx"
OUTPUT_DIR = ROOT / "generated"
ASSET_DIR = OUTPUT_DIR / "saheli_assets"
OUTPUT_PPTX = OUTPUT_DIR / "Saheli_WitchHunt_2026_Deck.pptx"

FONT_BOLD = "/usr/share/fonts/liberation/LiberationSans-Bold.ttf"
FONT_REGULAR = "/usr/share/fonts/liberation/LiberationSans-Regular.ttf"


WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(18, 18, 18)
GREEN = RGBColor(34, 79, 62)
LIGHT_GREEN = RGBColor(111, 141, 104)
GOLD = RGBColor(204, 170, 88)
PURPLE = RGBColor(84, 47, 110)
SOFT_WHITE = RGBColor(244, 240, 233)
MUTED = RGBColor(198, 203, 196)


def set_shape_text(
    shape,
    text: str,
    *,
    font_size: float | None = None,
    color: RGBColor | None = None,
    bold: bool | None = None,
    align: PP_ALIGN | None = None,
    valign: MSO_ANCHOR | None = None,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    if valign is not None:
        tf.vertical_anchor = valign

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            if color is not None:
                run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_text(
        shape,
        text,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        valign=MSO_ANCHOR.TOP,
    )
    return shape


def add_round_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: RGBColor,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    return shape


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def replace_picture(slide, picture_shape, image_path: Path):
    left, top, width, height = picture_shape.left, picture_shape.top, picture_shape.width, picture_shape.height
    remove_shape(picture_shape)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def make_persona_card(path: Path, title: str, subtitle: str, accent_hex: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    width, height = 900, 1100
    img = Image.new("RGB", (width, height), "#f3efe7")
    draw = ImageDraw.Draw(img)

    accent = tuple(int(accent_hex[i : i + 2], 16) for i in (0, 2, 4))
    green = (37, 74, 59)
    gold = (204, 170, 88)

    draw.rounded_rectangle((0, 0, width, height), radius=56, fill=(243, 239, 231))
    draw.rectangle((0, 0, int(width * 0.23), height), fill=green)
    draw.ellipse((120, 110, 690, 680), fill=accent)
    draw.ellipse((255, 205, 555, 505), fill=(246, 227, 200))
    draw.rounded_rectangle((205, 445, 610, 850), radius=160, fill=(246, 227, 200))
    draw.ellipse((110, 70, 730, 640), outline=gold, width=8)
    draw.rounded_rectangle((90, 905, 720, 1010), radius=26, fill=green)

    title_font = ImageFont.truetype(FONT_BOLD, 88)
    subtitle_font = ImageFont.truetype(FONT_REGULAR, 42)
    meta_font = ImageFont.truetype(FONT_BOLD, 36)

    draw.text((110, 918), title.upper(), fill=(255, 255, 255), font=title_font)
    draw.text((112, 1018), subtitle, fill=green, font=subtitle_font)
    draw.text((60, 50), "SAHELI", fill=(255, 255, 255), font=meta_font)

    img.save(path)


def update_slide_1(prs: Presentation):
    slide = prs.slides[0]
    set_shape_text(slide.shapes[2], "Proposal by\nSAHELI", font_size=28, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[3],
        "Voice-native, stealth-first support for women on shared phones.\n"
        "Evidence vault, rights guidance, schemes, health screening, and SOS.\n"
        "WitchHunt 2026 | HopeWorks x AI4India x Skills Cafe",
        font_size=12,
        color=SOFT_WHITE,
    )


def update_slide_2(prs: Presentation):
    slide = prs.slides[1]
    set_shape_text(slide.shapes[4], "CRISIS IN NUMBERS", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    metrics = [
        ("43%", "women do not own a smartphone"),
        ("58%", "rural women use shared devices"),
        ("67%", "harassment cases go unreported"),
        ("26%", "DV conviction rate due to evidence gaps"),
        ("74%", "Nirbhaya Fund still sits unspent"),
    ]
    top_boxes = [3, 9, 10, 11, 12]
    bottom_boxes = [2, 13, 14, 15, 16]
    card_shapes = [0, 5, 6, 7, 8]
    for (number, label), top_idx, bottom_idx, card_idx in zip(metrics, top_boxes, bottom_boxes, card_shapes, strict=True):
        set_shape_text(slide.shapes[top_idx], "", font_size=1)
        set_shape_text(slide.shapes[bottom_idx], label, font_size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        card = slide.shapes[card_idx]
        add_textbox(
            slide,
            card.left.inches + 0.05,
            card.top.inches + 0.42,
            card.width.inches - 0.10,
            0.5,
            number,
            font_size=26,
            color=GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def update_slide_3(prs: Presentation):
    slide = prs.slides[2]
    set_shape_text(slide.shapes[5], "Problem Statement", font_size=18, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[6],
        "Shared phones make safety and rights invisible",
        font_size=15,
        color=WHITE,
        bold=True,
    )
    set_shape_text(slide.shapes[8], "Problem Statement Expansion", font_size=12, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[9],
        "Most women-focused safety tools assume private smartphones, stable internet, and high digital literacy. "
        "Saheli starts from the actual constraint set: monitored devices, Hindi-first voice interaction, intermittent "
        "connectivity, and the need to document abuse without leaving traces.",
        font_size=10.5,
        color=SOFT_WHITE,
    )
    set_shape_text(slide.shapes[10], "How it impacts the user or client", font_size=11, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[11],
        "Abuse stays undocumented and cases weaken.\nWomen miss legal aid, schemes, and care.\n"
        "Institutions exist, but the last-mile access layer fails.",
        font_size=10.5,
        color=SOFT_WHITE,
    )
    set_shape_text(slide.shapes[13], "How it manifests for the user or client", font_size=11, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[14],
        "Safety apps are visible on shared devices.\nComplex UI excludes WhatsApp-first users.\n"
        "SOS-only tools ignore legal literacy, documentation, and recovery.",
        font_size=10.5,
        color=SOFT_WHITE,
    )
    remove_shape(slide.shapes[16])
    remove_shape(slide.shapes[15])
    remove_shape(slide.shapes[12])
    add_textbox(slide, 3.95, 4.48, 4.6, 0.22, "Data anchors: NSSO, IAMAI, GSMA, NCRB, NARI, NALSA", font_size=8.5, color=MUTED)


def update_slide_4(prs: Presentation):
    slide = prs.slides[3]
    set_shape_text(
        slide.shapes[3],
        "Saheli is a stealth-first, voice-native companion that turns one sentence into action: capture evidence, "
        "understand rights, check scheme eligibility, log health risk, or trigger help. It works across app, "
        "offline PWA, and SMS flows so support still reaches women on shared or low-end phones.",
        font_size=12,
        color=WHITE,
    )
    set_shape_text(slide.shapes[4], "Project Details", font_size=24, color=WHITE, bold=True)
    set_shape_text(slide.shapes[5], "SAHELI", font_size=24, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[6],
        "Women 18-45 in tier-2/3 India: shared-phone households, rural users,\nlow-literacy users, and feature-phone users.",
        font_size=10.5,
        color=WHITE,
    )


def update_slide_5(prs: Presentation):
    slide = prs.slides[4]
    set_shape_text(slide.shapes[3], "Saheli Design Principles", font_size=22, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[1],
        "Voice-first: conversation instead of navigation\n"
        "Stealth-first: hidden inside a devotional cover app\n"
        "Offline-first: SMS, cached PWA, and app tiers\n"
        "Evidence-first: encrypted, hashed, timestamped records\n"
        "Hindi-first: plain-language rights, health, and scheme support\n"
        "Real-world distribution: OSCs, ASHAs, SHGs, and legal aid",
        font_size=13,
        color=WHITE,
    )


def update_persona_slide(
    slide,
    *,
    banner: str,
    name: str,
    facts: str,
    bio: str,
    traits: str,
    triggers: str,
    problems: str,
    goals: str,
    image_path: Path,
):
    facts_shape = slide.shapes[1]
    bio_shape = slide.shapes[2]
    traits_heading = slide.shapes[4]
    traits_shape = slide.shapes[5]
    trigger_heading = slide.shapes[7]
    trigger_shape = slide.shapes[8]
    problems_heading = slide.shapes[11]
    problems_shape = slide.shapes[12]
    goals_heading = slide.shapes[14]
    goals_shape = slide.shapes[15]
    name_shape = slide.shapes[17].shapes[1]
    picture_shape = slide.shapes[18]
    banner_shape = slide.shapes[19]

    set_shape_text(facts_shape, facts, font_size=12.5, color=WHITE, bold=True)
    set_shape_text(bio_shape, bio, font_size=10.5, color=BLACK)
    set_shape_text(traits_heading, "ABOUT THE USER", font_size=11, color=GREEN, bold=True)
    set_shape_text(traits_shape, traits, font_size=10, color=GREEN)
    set_shape_text(trigger_heading, "TRIGGER MOMENTS", font_size=11, color=GREEN, bold=True)
    set_shape_text(trigger_shape, triggers, font_size=10, color=GREEN)
    set_shape_text(problems_heading, "PROBLEMS", font_size=11, color=GREEN, bold=True)
    set_shape_text(problems_shape, problems, font_size=10, color=GREEN)
    set_shape_text(goals_heading, "GOALS AND NEEDS", font_size=11, color=GREEN, bold=True)
    set_shape_text(goals_shape, goals, font_size=10, color=GREEN)
    set_shape_text(name_shape, name.upper(), font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    set_shape_text(banner_shape, banner, font_size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    remove_shape(slide.shapes[21])
    replace_picture(slide, picture_shape, image_path)


def update_slide_6_to_8(prs: Presentation):
    priya_card = ASSET_DIR / "priya.png"
    sunita_card = ASSET_DIR / "sunita.png"
    meena_card = ASSET_DIR / "meena.png"
    make_persona_card(priya_card, "Priya", "Lucknow | shared phone", "9D7B4F")
    make_persona_card(sunita_card, "Sunita", "Nashik | basic smartphone", "7E9C6A")
    make_persona_card(meena_card, "Meena", "Varanasi | feature phone", "8F5F9D")

    update_persona_slide(
        prs.slides[5],
        banner="USER PERSONA 1",
        name="Priya",
        facts="35 years old\nLucknow\nShared-phone mother\n2 children",
        bio=(
            "Priya lives with an abusive husband who controls the phone, the money, and what apps can stay installed. "
            "She needs a way to document violence, understand her rights, and reach support without exposing herself."
        ),
        traits="Cautious\nResourceful\nProtective\nHindi-first\nPrivacy-sensitive",
        triggers="Phone use is monitored\nNeeds proof without a visible app\nWants rights explained safely",
        problems="No secure place to store evidence\nLimited knowledge of DV protections\nCannot openly ask for help",
        goals="Capture proof invisibly\nReach OSC or trusted contacts\nUnderstand maintenance, custody, and property",
        image_path=priya_card,
    )
    update_persona_slide(
        prs.slides[6],
        banner="USER PERSONA 2",
        name="Sunita",
        facts="29 years old\nNashik\nTeacher\nBasic smartphone",
        bio=(
            "Sunita has lived with severe period pain for years and has learned to call it normal because care is costly, "
            "confusing, and never structured around her daily reality."
        ),
        traits="Resilient\nUnderdiagnosed\nBudget-conscious\nWorking daily\nSeeks relief",
        triggers="Pain is normalized at home\nDoctor visits feel expensive\nNeeds simple Hindi guidance",
        problems="No symptom history to share\nNo structured screening flow\nAffordable care is hard to find",
        goals="Track symptoms over time\nReach low-cost medicine and care\nCarry a doctor-ready summary",
        image_path=sunita_card,
    )
    update_persona_slide(
        prs.slides[7],
        banner="USER PERSONA 3",
        name="Meena",
        facts="42 years old\nVaranasi\nFeature phone user\nNo internet access",
        bio=(
            "Meena does not have app-install capacity, data access, or a habit of digital healthcare. If support depends "
            "on a modern smartphone, she is excluded before the journey even starts."
        ),
        traits="Low-literacy\nRural\nPersistent\nNever screened\nTrusts ASHA workers",
        triggers="No data plan or app store access\nNeeds SMS-first guidance\nLives far from formal care",
        problems="No installable app path\nScreening information never reaches her\nHigh-risk signs go unchecked",
        goals="Screen through SMS or USSD\nGet nearest PHC guidance\nReceive follow-up through ASHA",
        image_path=meena_card,
    )


def update_slide_9(prs: Presentation):
    slide = prs.slides[8]
    set_shape_text(slide.shapes[0], "Our Solution", font_size=24, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[1],
        "A stealth, voice-native rights and evidence companion for women on shared phones",
        font_size=13,
        color=SOFT_WHITE,
    )
    set_shape_text(
        slide.shapes[8].shapes[2],
        "What we build?\nStealth app + offline PWA + SMS layer\n\n"
        "Who is it for?\nWomen on shared or low-end phones in tier-2/3 India\n\n"
        "What problem does it address?\nSafe access to rights, evidence, schemes, and care\n\n"
        "How would it work in the real world?\nVoice triggers workflows; quick exit, cached data, and helpline routing keep it usable under pressure\n\n"
        "What impact would it create?\nMoves users from awareness to documented action\n\n"
        "Why is it worth building?\nBecause current safety products fail exactly where privacy and accessibility matter most",
        font_size=10.5,
        color=WHITE,
    )


def update_slide_10(prs: Presentation):
    slide = prs.slides[9]
    set_shape_text(slide.shapes[0], "Our Solution Continue", font_size=24, color=WHITE, bold=True)
    add_textbox(
        slide,
        4.15,
        0.95,
        5.0,
        0.55,
        "Five capabilities, one conversation",
        font_size=18,
        color=WHITE,
        bold=True,
    )
    items = [
        ("01", "Evidence Vault", "Hashed photos, recordings, documents, timestamps, GPS"),
        ("02", "Know Your Rights", "PWDVA, POSH, maintenance, property, and legal aid workflows"),
        ("03", "Safety Net", "Stealth SOS, offline SMS fallback, cached helplines, OSC locator"),
        ("04", "Entitlements", "OCR-based eligibility checks for government and NGO schemes"),
        ("05", "Health Journal", "Voice screening for EPDS, PCOS, cervical risk, plus doctor card"),
    ]
    top = 1.65
    for number, heading, body in items:
        add_round_box(slide, 4.1, top, 4.75, 0.72, RGBColor(25, 25, 25), LIGHT_GREEN)
        add_textbox(slide, 4.28, top + 0.09, 0.45, 0.3, number, font_size=14, color=GOLD, bold=True)
        add_textbox(slide, 4.75, top + 0.06, 1.95, 0.28, heading, font_size=14, color=WHITE, bold=True)
        add_textbox(slide, 4.75, top + 0.31, 3.75, 0.3, body, font_size=9.5, color=SOFT_WHITE)
        top += 0.78


def update_slide_11(prs: Presentation):
    slide = prs.slides[10]
    set_shape_text(slide.shapes[4], "Artificial Intelligence", font_size=26, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[5],
        "AI is the interface layer, not the product. It converts complex legal, health, and scheme workflows into safe, "
        "spoken guidance on low-end Indian devices.",
        font_size=12,
        color=SOFT_WHITE,
    )
    set_shape_text(slide.shapes[6], "USE OF", font_size=14, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[7],
        "",
        font_size=1,
    )
    remove_shape(slide.shapes[7])
    add_round_box(slide, 1.45, 2.0, 5.3, 1.25, RGBColor(32, 46, 87), RGBColor(122, 145, 209))
    add_textbox(
        slide,
        1.65,
        2.12,
        4.95,
        1.0,
        "Saaras V3 STT: Hindi-first voice capture across 22+ languages\n"
        "Bulbul V3 TTS: empathetic spoken guidance for low-literacy users\n"
        "Sarvam-30B + legal RAG: verified rights and entitlement workflows\n"
        "Edge + Vision + Samvaad: offline AI, OCR of documents, and omnichannel delivery",
        font_size=10.5,
        color=WHITE,
    )


def update_slide_12(prs: Presentation):
    slide = prs.slides[11]
    set_shape_text(slide.shapes[2], "Challenges and Risks", font_size=26, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[3],
        "Trust and safety:\nAI must never sound like legal or medical advice; guardrails and escalation are required.",
        font_size=12,
        color=WHITE,
    )
    set_shape_text(
        slide.shapes[4],
        "Stealth failure:\nAny visible trace can endanger the user; quick exit, hidden notifications, and cover mode are non-negotiable.",
        font_size=12,
        color=WHITE,
    )
    set_shape_text(
        slide.shapes[5],
        "Localization drift:\nDialect accuracy and cached knowledge can age; curated updates and human review keep the system reliable.",
        font_size=12,
        color=WHITE,
    )
    set_shape_text(
        slide.shapes[8],
        "",
        font_size=1,
    )
    remove_shape(slide.shapes[8])
    add_round_box(slide, 1.35, 1.9, 4.95, 0.95, RGBColor(31, 40, 73), RGBColor(99, 117, 177))
    add_textbox(
        slide,
        1.58,
        2.1,
        4.5,
        0.55,
        "The product only works if privacy, accuracy, and last-mile delivery outperform the average safety app.",
        font_size=12,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def update_slide_13(prs: Presentation):
    slide = prs.slides[12]
    set_shape_text(slide.shapes[0], "Technical Request for Facilitated Datasets", font_size=20, color=WHITE, bold=True)
    set_shape_text(slide.shapes[1], "Specification of Required Data", font_size=12, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[2],
        "Data Type: Indic speech from women in low-resource dialects; verified service directories for OSCs, PHCs, Jan Aushadhi, and legal aid; anonymized women's health screening data.\n\n"
        "Volume & Depth: District-level coverage for high-need states plus enough speech hours to validate rural Hindi and code-mixed usage.\n\n"
        "Granularity: Village or block-level service routing where available; speaker metadata by dialect band.",
        font_size=9.5,
        color=WHITE,
    )
    set_shape_text(slide.shapes[3], 'The "Ask" vs. "Leveraged" Mapping', font_size=12, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[4],
        "What we have: Bare Acts, NALSA resources, WHO and UNICEF protocols, public scheme data, and open Indic voice baselines.\n\n"
        "What we need: Facilitated speech, service-location, and validation datasets that make the product accurate for underserved women in real contexts.",
        font_size=9.5,
        color=WHITE,
    )
    set_shape_text(slide.shapes[5], "Technical Justification", font_size=12, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[6],
        "Accuracy: improves STT for rural voices and code-mixed queries.\n\n"
        "Validation: grounds EPDS, PCOS, and cervical-risk flows in real public-health distributions.\n\n"
        "Routing: makes offline referral and entitlement matching trustworthy at district scale.",
        font_size=9.5,
        color=WHITE,
    )
    set_shape_text(slide.shapes[8], "Environment & Security Requirements", font_size=12, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[9],
        "Access Method: secure API, sandbox, or encrypted batch download.\n\n"
        "Security: minimum-retention storage, access controls, encryption at rest, and DPDP-aligned handling for all sensitive fields.",
        font_size=9.5,
        color=WHITE,
    )
    set_shape_text(slide.shapes[11], "Request focus: voice accuracy, service routing, and health validation.", font_size=10, color=MUTED)


def update_slide_14(prs: Presentation):
    slide = prs.slides[13]
    set_shape_text(slide.shapes[0], "Anticipated Results", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    set_shape_text(
        slide.shapes[1],
        "Year 1 targets anchored to documented action, not vanity metrics.",
        font_size=11,
        color=SOFT_WHITE,
    )
    set_shape_text(slide.shapes[3], "Saheli Impact Model", font_size=16, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[4].shapes[4],
        "10,000+ women document evidence safely; Priya-like users benefit first where proof is the biggest missing layer.",
        font_size=11,
        color=WHITE,
    )
    set_shape_text(
        slide.shapes[4].shapes[5],
        "50,000+ complete the rights module; 25,000+ run scheme eligibility checks; 20,000+ interact through SMS-first flows.",
        font_size=11,
        color=WHITE,
    )
    set_shape_text(
        slide.shapes[4].shapes[6],
        "100,000+ voice health screenings, distributed through ASHAs, SHGs, OSCs, and lightweight app/PWA/SMS infrastructure.",
        font_size=11,
        color=WHITE,
    )
    set_shape_text(slide.shapes[5].shapes[1], "TARGETS", font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def update_slide_15(prs: Presentation):
    slide = prs.slides[14]
    set_shape_text(slide.shapes[1], "THANK YOU", font_size=34, color=WHITE, bold=True)
    set_shape_text(
        slide.shapes[2],
        "Saheli | one conversation for every right, every scheme, and every tool",
        font_size=16,
        color=SOFT_WHITE,
    )


def build_deck():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(TEMPLATE))
    update_slide_1(prs)
    update_slide_2(prs)
    update_slide_3(prs)
    update_slide_4(prs)
    update_slide_5(prs)
    update_slide_6_to_8(prs)
    update_slide_9(prs)
    update_slide_10(prs)
    update_slide_11(prs)
    update_slide_12(prs)
    update_slide_13(prs)
    update_slide_14(prs)
    update_slide_15(prs)
    prs.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    build_deck()
