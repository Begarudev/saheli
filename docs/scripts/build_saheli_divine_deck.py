from __future__ import annotations

from pathlib import Path
import math
import random

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "generated"
ASSET_DIR = OUT_DIR / "divine_assets"
OUT_PPTX = OUT_DIR / "Saheli_WitchHunt_2026_Divine_Deck.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_HEAD = "Noto Serif"
FONT_BODY = "Noto Sans"
FONT_DEV = "Noto Sans Devanagari"

PIL_HEAD_BOLD = "/usr/share/fonts/noto/NotoSerif-Bold.ttf"
PIL_HEAD = "/usr/share/fonts/noto/NotoSerif-Regular.ttf"
PIL_BODY = "/usr/share/fonts/noto/NotoSans-Regular.ttf"
PIL_BODY_BOLD = "/usr/share/fonts/noto/NotoSans-Bold.ttf"
PIL_DEV = "/usr/share/fonts/noto/NotoSansDevanagari-Bold.ttf"

INDIGO = RGBColor(24, 22, 42)
PLUM = RGBColor(70, 52, 88)
LOTUS = RGBColor(190, 122, 136)
SINDOOR = RGBColor(177, 82, 75)
HALDI = RGBColor(210, 163, 95)
PEACOCK = RGBColor(62, 108, 103)
PAPER = RGBColor(246, 238, 228)
SANDAL = RGBColor(234, 219, 195)
INK = RGBColor(35, 29, 32)
SMOKE = RGBColor(93, 84, 88)
WHITE = RGBColor(249, 247, 243)
MIST = RGBColor(221, 214, 205)


def rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def add_noise(img: Image.Image, alpha: int, count: int = 3000):
    px = img.load()
    w, h = img.size
    for _ in range(count):
        x = random.randint(0, w - 1)
        y = random.randint(0, h - 1)
        v = random.randint(180, 255)
        px[x, y] = (v, v, v, alpha)


def draw_mandala(draw: ImageDraw.ImageDraw, cx: int, cy: int, radius: int, color: tuple[int, int, int, int]):
    for i in range(12):
        angle = math.radians(i * 30)
        x1 = cx + int(math.cos(angle) * radius * 0.4)
        y1 = cy + int(math.sin(angle) * radius * 0.4)
        x2 = cx + int(math.cos(angle) * radius)
        y2 = cy + int(math.sin(angle) * radius)
        draw.line((x1, y1, x2, y2), fill=color, width=2)
        draw.ellipse((x2 - 7, y2 - 7, x2 + 7, y2 + 7), outline=color, width=2)
    draw.ellipse((cx - radius // 2, cy - radius // 2, cx + radius // 2, cy + radius // 2), outline=color, width=3)
    draw.ellipse((cx - radius // 4, cy - radius // 4, cx + radius // 4, cy + radius // 4), outline=color, width=2)


def draw_lotus(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float, color: tuple[int, int, int, int]):
    petals = [
        [(0, 40), (35, -10), (70, 40)],
        [(35, 35), (70, -15), (105, 35)],
        [(70, 40), (105, -10), (140, 40)],
        [(20, 58), (55, 15), (90, 58)],
        [(50, 60), (85, 18), (120, 60)],
    ]
    for petal in petals:
        pts = [(x + int(px * scale), y + int(py * scale)) for px, py in petal]
        draw.polygon(pts, outline=color)
    draw.arc((x + int(5 * scale), y + int(45 * scale), x + int(135 * scale), y + int(95 * scale)), 180, 360, fill=color, width=2)


def make_bg_cover(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#17162a"))
    ol = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(ol)
    d.ellipse((1080, 50, 1560, 530), fill=(208, 163, 95, 48))
    d.ellipse((1160, 130, 1480, 450), outline=(241, 217, 175, 150), width=2)
    d.ellipse((1130, 100, 1510, 480), outline=(188, 122, 136, 90), width=4)
    for offset in range(4):
        draw_lotus(d, 1120 + offset * 40, 520 + offset * 12, 1.4, (238, 215, 181, 110))
    draw_mandala(d, 1360, 200, 92, (232, 214, 184, 70))
    for y in range(140, 780, 34):
        d.ellipse((980, y, 985, y + 5), fill=(238, 215, 181, 90))
    d.arc((980, 100, 1510, 660), 210, 330, fill=(99, 140, 129, 140), width=3)
    ol = ol.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, ol)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_noise(grain, 12)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_bg_paper(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#f6eee4"))
    ol = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(ol)
    d.ellipse((1040, -120, 1610, 340), fill=(221, 194, 150, 52))
    d.ellipse((-180, 620, 350, 980), fill=(111, 149, 128, 32))
    draw_mandala(d, 1320, 170, 84, (194, 150, 91, 70))
    draw_mandala(d, 1460, 750, 64, (101, 131, 116, 50))
    d.arc((1020, 40, 1540, 560), 200, 315, fill=(112, 151, 129, 95), width=3)
    d.arc((1070, 90, 1490, 510), 200, 315, fill=(192, 132, 142, 60), width=2)
    ol = ol.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, ol)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_noise(grain, 8, 4200)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_bg_deep(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#1b1631"))
    ol = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(ol)
    d.ellipse((980, 120, 1540, 680), fill=(71, 52, 88, 70))
    d.ellipse((-220, 540, 360, 980), fill=(75, 118, 108, 42))
    for i in range(5):
        draw_lotus(d, 1020 + i * 55, 590 - i * 18, 1.1, (223, 190, 140, 95))
    draw_mandala(d, 1260, 170, 70, (225, 205, 173, 60))
    d.arc((930, 40, 1540, 640), 200, 330, fill=(214, 170, 97, 110), width=4)
    d.arc((960, 70, 1510, 610), 200, 330, fill=(111, 149, 128, 90), width=2)
    ol = ol.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, ol)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_noise(grain, 10)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def make_bg_rose(path: Path):
    img = Image.new("RGBA", (1600, 900), rgb("#f4e7e3"))
    ol = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(ol)
    d.ellipse((1120, -60, 1660, 420), fill=(193, 120, 136, 54))
    d.ellipse((-220, 620, 310, 960), fill=(206, 167, 98, 30))
    draw_mandala(d, 1310, 170, 82, (175, 94, 82, 60))
    draw_mandala(d, 1500, 740, 56, (103, 137, 123, 50))
    d.arc((1060, 80, 1530, 550), 190, 320, fill=(102, 138, 124, 70), width=2)
    ol = ol.filter(ImageFilter.GaussianBlur(10))
    img = Image.alpha_composite(img, ol)
    grain = Image.new("RGBA", img.size, (0, 0, 0, 0))
    add_noise(grain, 7, 4000)
    img = Image.alpha_composite(img, grain)
    img.save(path)


def build_backgrounds():
    ensure_dirs()
    make_bg_cover(ASSET_DIR / "bg_cover.png")
    make_bg_paper(ASSET_DIR / "bg_paper.png")
    make_bg_deep(ASSET_DIR / "bg_deep.png")
    make_bg_rose(ASSET_DIR / "bg_rose.png")


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def new_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, filename: str):
    slide.shapes.add_picture(str(ASSET_DIR / filename), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def style_run(run, font_name: str, size: float, color: RGBColor, bold: bool = False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, x, y, w, h, text, *, font_name=FONT_BODY, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    style_run(r, font_name, size, color, bold)
    return box


def add_bullets(slide, x, y, w, h, items, *, size=13, color=WHITE, bullet=HALDI, font_name=FONT_BODY, spacing=1.18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• "
        style_run(r1, font_name, size, bullet, True)
        r2 = p.add_run()
        r2.text = item
        style_run(r2, font_name, size, color, False)
    return box


def add_card(slide, x, y, w, h, *, fill: RGBColor, line: RGBColor | None = None, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    shp.line.width = Pt(1.2 if line else 0)
    return shp


def add_chip(slide, x, y, w, h, label, *, fill, text_color, line=None, size=10.5):
    add_card(slide, x, y, w, h, fill=fill, line=line)
    add_text(slide, x, y + 0.04, w, h - 0.04, label, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def add_rule(slide, x, y, w, color):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    return shp


def slide_cover(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_cover.png")
    add_text(s, 0.85, 0.6, 2.6, 0.25, "WITCHHUNT 2026", size=11, color=MIST, bold=True)
    add_text(s, 0.85, 1.1, 4.2, 0.85, "SAHELI", font_name=FONT_HEAD, size=31, color=WHITE, bold=True)
    add_text(s, 0.88, 1.9, 2.1, 0.35, "सहेली", font_name=FONT_DEV, size=21, color=RGBColor(235, 222, 206), bold=True)
    add_text(s, 0.85, 2.45, 5.2, 1.0, "A voice-native, stealth-first companion for rights, evidence, safety, health, and dignity.", size=18, color=WHITE)
    add_card(s, 0.85, 4.15, 4.85, 1.5, fill=RGBColor(40, 32, 55), line=RGBColor(107, 90, 122))
    add_text(s, 1.12, 4.45, 4.3, 0.78, "Rooted in the realities of shared phones, low digital literacy, and Hindi-first support in India.", size=16, color=WHITE)
    chips = [
        ("Evidence Vault", 0.85, 6.23, 1.8),
        ("Rights Guidance", 2.85, 6.23, 1.9),
        ("Stealth SOS", 4.95, 6.23, 1.55),
        ("Health Journal", 6.68, 6.23, 1.75),
        ("Scheme Access", 8.63, 6.23, 1.75),
    ]
    for label, x, y, w in chips:
        add_chip(s, x, y, w, 0.46, label, fill=RGBColor(35, 30, 46), text_color=WHITE, line=RGBColor(100, 88, 117), size=9.6)
    add_text(s, 9.9, 5.95, 2.3, 0.4, "HopeWorks x AI4India x Skills Cafe", size=11, color=MIST, align=PP_ALIGN.RIGHT)


def slide_insight(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.8, 0.55, 2.5, 0.25, "THE ACCESS CRISIS", size=11, color=SMOKE, bold=True)
    add_text(s, 0.8, 0.95, 5.5, 0.9, "Women’s safety is still being designed for the wrong device, the wrong interface, and the wrong conditions.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    stats = [
        ("43%", "women without a personal smartphone"),
        ("58%", "rural women on shared devices"),
        ("67%", "harassment that still goes unreported"),
        ("26%", "DV conviction rate constrained by proof gaps"),
    ]
    positions = [(0.9, 2.15), (2.95, 2.15), (0.9, 4.0), (2.95, 4.0)]
    for (num, txt), (x, y) in zip(stats, positions, strict=True):
        add_card(s, x, y, 1.75, 1.45, fill=WHITE, line=SANDAL)
        add_text(s, x + 0.12, y + 0.15, 1.4, 0.35, num, font_name=FONT_HEAD, size=22, color=PEACOCK, bold=True)
        add_text(s, x + 0.12, y + 0.68, 1.42, 0.55, txt, size=10, color=SMOKE)
    add_card(s, 6.15, 1.9, 6.1, 4.9, fill=RGBColor(33, 28, 44), line=RGBColor(103, 88, 118))
    add_text(s, 6.5, 2.2, 1.8, 0.22, "CORE INSIGHT", size=11, color=RGBColor(235, 210, 174), bold=True)
    add_text(s, 6.5, 2.55, 5.2, 1.5, "The laws exist. The schemes exist. The public infrastructure exists.\n\nWhat is missing is a private, language-first, offline-tolerant path that women can actually use when their phone is not fully theirs.", size=17, color=WHITE)
    add_bullets(s, 6.5, 4.55, 5.0, 1.35, [
        "Shared devices turn visible safety apps into liabilities.",
        "Complex English-heavy UI excludes WhatsApp-first users.",
        "Emergency-only tools ignore proof, rights, and recovery.",
    ], size=12, color=MIST, bullet=HALDI)
    add_text(s, 0.82, 6.88, 4.8, 0.18, "NSSO, IAMAI, GSMA, NCRB, NARI, NALSA", size=9.3, color=SMOKE)


def slide_failure_modes(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_deep.png")
    add_text(s, 0.8, 0.55, 2.8, 0.25, "WHERE CURRENT TOOLS BREAK", size=11, color=MIST, bold=True)
    add_text(s, 0.8, 0.95, 5.5, 0.86, "The mainstream women’s-safety product stack still fails at the exact moment friction becomes risk.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    cards = [
        ("Private phone assumption", "A visible safety app can itself create danger on a monitored device."),
        ("Internet dependence", "GPS and chatbot-heavy flows collapse when connectivity is weak."),
        ("Emergency-only framing", "SOS without proof, rights, or recovery leaves women underprepared."),
        ("Weak institutional routing", "Trust remains low when the response loop feels opaque or hostile."),
        ("High cognitive load", "If it feels harder than a voice note, many users will never reach the useful part."),
    ]
    coords = [(0.95, 2.25), (4.2, 2.25), (7.45, 2.25), (2.57, 4.65), (5.82, 4.65)]
    for idx, ((title, body), (x, y)) in enumerate(zip(cards, coords, strict=True), start=1):
        add_card(s, x, y, 2.85, 1.72, fill=RGBColor(39, 31, 53), line=RGBColor(105, 90, 121))
        add_text(s, x + 0.18, y + 0.18, 0.5, 0.22, f"{idx:02d}", size=11, color=HALDI, bold=True)
        add_text(s, x + 0.18, y + 0.48, 2.2, 0.4, title, size=13.2, color=WHITE, bold=True)
        add_text(s, x + 0.18, y + 0.94, 2.28, 0.5, body, size=10.2, color=MIST)
    add_card(s, 0.95, 6.55, 11.45, 0.46, fill=RGBColor(35, 29, 47), line=RGBColor(100, 88, 117))
    add_text(s, 1.12, 6.64, 11.0, 0.16, "Saheli begins with conversation, camouflage, and culturally familiar interaction patterns, not with panic-button assumptions.", size=11.1, color=WHITE, align=PP_ALIGN.CENTER)


def slide_what_is(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_rose.png")
    add_text(s, 0.82, 0.55, 2.0, 0.25, "WHAT SAHELI IS", size=11, color=SMOKE, bold=True)
    add_text(s, 0.82, 0.95, 5.0, 0.82, "A feminine, invisible layer of support inside an ordinary-looking app shell.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    add_card(s, 0.95, 1.95, 4.8, 4.85, fill=RGBColor(37, 32, 49), line=RGBColor(109, 92, 123))
    add_text(s, 1.22, 2.18, 2.0, 0.2, "ONE-LINE PITCH", size=11, color=RGBColor(235, 210, 174), bold=True)
    add_text(s, 1.22, 2.55, 4.0, 1.6, "Saheli is a voice-native companion hidden behind a devotional cover app.\n\nIt gives women access to an encrypted evidence vault, rights guidance, scheme support, health screening, and stealth SOS in Hindi, even when the device is shared.", size=16.5, color=WHITE)
    add_text(s, 1.22, 5.2, 4.0, 0.62, "Built for women in tier-2/3 India, including shared-phone households and feature-phone pathways.", size=11.2, color=MIST)
    pillars = [
        ("Voice-native", "Conversation replaces menu-driven friction."),
        ("Stealth-first", "Hidden entry, quick exit, devotional cover."),
        ("Offline-first", "SMS and cached PWA support the low-connectivity edge."),
        ("Evidence-first", "Proof is encrypted, hashed, and preserved from the start."),
    ]
    coords = [(6.3, 2.05), (9.35, 2.05), (6.3, 4.5), (9.35, 4.5)]
    for (title, body), (x, y) in zip(pillars, coords, strict=True):
        add_card(s, x, y, 2.7, 1.85, fill=WHITE, line=SANDAL)
        add_text(s, x + 0.18, y + 0.2, 2.1, 0.32, title, size=13, color=PEACOCK, bold=True)
        add_text(s, x + 0.18, y + 0.62, 2.12, 0.68, body, size=10.3, color=SMOKE)


def slide_capabilities(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_cover.png")
    add_text(s, 0.82, 0.55, 3.5, 0.25, "ONE CONVERSATION, FIVE CAPABILITIES", size=11, color=MIST, bold=True)
    add_text(s, 0.82, 0.95, 5.4, 0.82, "She speaks once. The product quietly routes her to the right support ritual.", font_name=FONT_HEAD, size=24, color=WHITE, bold=True)
    labels = [
        ("Evidence Vault", "Saboot rakhna hai", "Encrypted photos, recordings, threats, documents, timestamps."),
        ("Know Your Rights", "Mera kya haq hai?", "PWDVA, POSH, maintenance, custody, legal aid, and property workflows."),
        ("Safety Net", "Mujhe madad chahiye", "Stealth SOS, offline fallback, cached helplines, OSC routing."),
        ("Entitlements", "Sarkari yojana", "Eligibility checks for welfare, loans, and public support."),
        ("Health Journal", "Tabiyat theek nahi", "Voice screening, symptom tracking, and doctor-ready summaries."),
    ]
    y = 1.9
    for title, trigger, body in labels:
        add_card(s, 5.15, y, 7.0, 0.94, fill=RGBColor(42, 34, 57), line=RGBColor(112, 93, 125))
        add_text(s, 5.42, y + 0.14, 2.25, 0.22, title, size=13, color=WHITE, bold=True)
        add_text(s, 7.7, y + 0.15, 1.55, 0.18, trigger, size=10.4, color=HALDI, bold=True)
        add_text(s, 5.42, y + 0.47, 6.0, 0.22, body, size=10.1, color=MIST)
        y += 0.99
    add_card(s, 0.95, 2.1, 3.4, 3.95, fill=RGBColor(36, 30, 49), line=RGBColor(103, 89, 118))
    add_text(s, 1.25, 2.42, 2.8, 1.2, "The interaction stays familiar because it feels like a voice note, not like learning a new system in a moment of stress.", size=16, color=WHITE)
    add_bullets(s, 1.25, 4.35, 2.75, 1.2, [
        "Hindi-first phrasing",
        "No visible safety branding",
        "Works across app, PWA, and SMS",
    ], size=11.5, color=MIST, bullet=HALDI)


def slide_stealth_flow(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.82, 0.55, 2.6, 0.25, "STEALTH EXPERIENCE", size=11, color=SMOKE, bold=True)
    add_text(s, 0.82, 0.95, 5.6, 0.82, "Designed to feel culturally familiar, practically useful, and visually unthreatening.", font_name=FONT_HEAD, size=23, color=INK, bold=True)
    steps = [
        ("01", "Daily Mantras", "The cover app looks devotional, not defensive."),
        ("02", "Hidden unlock", "Five taps plus PIN reveals the real layer."),
        ("03", "Speak naturally", "Voice-first input lowers the threshold to ask for help."),
        ("04", "Route silently", "Rights, evidence, health, or SOS flows load beneath the shell."),
        ("05", "Exit without traces", "Quick exit and hidden notifications reduce exposure."),
    ]
    xs = [0.95, 3.42, 5.89, 8.36, 10.83]
    for (num, title, body), x in zip(steps, xs, strict=True):
        add_card(s, x, 3.0, 2.0, 2.45, fill=WHITE, line=SANDAL)
        add_text(s, x + 0.18, 3.2, 0.45, 0.2, num, size=11, color=SINDOOR, bold=True)
        add_text(s, x + 0.18, 3.55, 1.55, 0.42, title, size=12.7, color=PEACOCK, bold=True)
        add_text(s, x + 0.18, 4.12, 1.55, 0.76, body, size=9.8, color=SMOKE)
    for x in [2.78, 5.25, 7.72, 10.19]:
        add_rule(s, x, 4.18, 0.42, RGBColor(194, 164, 117))
    add_card(s, 0.95, 6.28, 11.4, 0.44, fill=RGBColor(34, 29, 46), line=RGBColor(104, 90, 118))
    add_text(s, 1.1, 6.36, 11.05, 0.16, "The product is protective because it is calm, ordinary-looking, and hard to weaponize against the user.", size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_personas(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_rose.png")
    add_text(s, 0.82, 0.55, 2.0, 0.25, "THREE WOMEN", size=11, color=SMOKE, bold=True)
    add_text(s, 0.82, 0.95, 5.8, 0.82, "Different lives. Different pressures. The same need for quiet, trustworthy support.", font_name=FONT_HEAD, size=23.5, color=INK, bold=True)
    personas = [
        ("Priya", "DV survivor on a shared phone", ["Needs covert proof capture", "Needs rights explained gently", "Needs a path to OSC and maintenance support"], PLUM),
        ("Sunita", "Teacher living with severe period pain", ["Needs symptom screening in Hindi", "Needs affordable care pathways", "Needs a doctor-ready summary"], PEACOCK),
        ("Meena", "Feature-phone user without internet", ["Needs an SMS-first route", "Needs referral without install friction", "Needs ASHA-linked follow-up"], RGBColor(143, 120, 74)),
    ]
    xs = [0.95, 4.25, 7.55]
    for (name, subtitle, bullets, accent), x in zip(personas, xs, strict=True):
        add_card(s, x, 2.05, 2.85, 4.75, fill=WHITE, line=SANDAL)
        add_rule(s, x, 2.05, 2.85, accent)
        add_text(s, x + 0.2, 2.34, 2.1, 0.3, name, font_name=FONT_HEAD, size=20, color=INK, bold=True)
        add_text(s, x + 0.2, 2.72, 2.2, 0.32, subtitle, size=10.8, color=SMOKE, bold=True)
        add_bullets(s, x + 0.2, 3.35, 2.25, 1.6, bullets, size=10.7, color=SMOKE, bullet=accent)
        add_chip(s, x + 0.2, 5.95, 0.88, 0.34, "USER", fill=PAPER, text_color=SMOKE, line=SANDAL, size=8.5)
        add_chip(s, x + 1.2, 5.95, 1.15, 0.34, "REAL NEED", fill=PAPER, text_color=SMOKE, line=SANDAL, size=8.5)


def slide_stack(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_deep.png")
    add_text(s, 0.82, 0.55, 3.0, 0.25, "AI STACK + PUBLIC INFRASTRUCTURE", size=11, color=MIST, bold=True)
    add_text(s, 0.82, 0.95, 5.8, 0.82, "Localized intelligence matters only when it can reach women through systems that already exist.", font_name=FONT_HEAD, size=23, color=WHITE, bold=True)
    add_card(s, 0.95, 1.95, 5.65, 5.0, fill=RGBColor(41, 33, 56), line=RGBColor(107, 90, 121))
    add_text(s, 1.22, 2.2, 1.8, 0.22, "AI LAYER", size=11, color=HALDI, bold=True)
    ai = [
        "Sarvam Saaras V3 for Indic speech-to-text",
        "Bulbul V3 for natural voice output",
        "Sarvam-30B plus domain-specific legal RAG",
        "Sarvam Edge for offline support",
        "Sarvam Vision for OCR on user documents",
        "Samvaad for omnichannel delivery across WhatsApp, web, and SMS",
    ]
    y = 2.62
    for item in ai:
        add_chip(s, 1.22, y, 4.95, 0.42, item, fill=RGBColor(52, 42, 67), text_color=WHITE, line=RGBColor(118, 101, 132), size=9.6)
        y += 0.54
    add_card(s, 6.9, 1.95, 5.45, 5.0, fill=RGBColor(36, 31, 50), line=RGBColor(105, 90, 121))
    add_text(s, 7.18, 2.2, 2.8, 0.22, "DISTRIBUTION LAYER", size=11, color=HALDI, bold=True)
    infra = [
        "733+ One Stop Centres",
        "9M Self-Help Groups",
        "1M+ ASHA workers",
        "District legal aid networks",
        "Jan Aushadhi and PHC touchpoints",
        "ABHA-linked digital health records",
    ]
    y = 2.62
    for item in infra:
        add_chip(s, 7.18, y, 4.65, 0.42, item, fill=RGBColor(48, 40, 63), text_color=WHITE, line=RGBColor(118, 101, 132), size=9.8)
        y += 0.54
    add_text(s, 7.18, 6.02, 4.45, 0.38, "Saheli is not replacing institutions. It is making them reachable through a more humane interface.", size=10.8, color=MIST)


def slide_impact(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_paper.png")
    add_text(s, 0.82, 0.55, 2.4, 0.25, "IMPACT + SCALE", size=11, color=SMOKE, bold=True)
    add_text(s, 0.82, 0.95, 5.5, 0.82, "Success should look like documented action, not vanity growth.", font_name=FONT_HEAD, size=24, color=INK, bold=True)
    stats = [
        ("10,000+", "women using the evidence vault"),
        ("50,000+", "rights journeys completed"),
        ("25,000+", "scheme eligibility checks"),
        ("100,000+", "voice-based health screenings"),
    ]
    xs = [0.95, 3.9, 6.85, 9.8]
    for (num, label), x in zip(stats, xs, strict=True):
        add_card(s, x, 2.2, 2.2, 1.55, fill=WHITE, line=SANDAL)
        add_text(s, x + 0.14, 2.38, 1.8, 0.32, num, font_name=FONT_HEAD, size=21, color=SINDOOR, bold=True)
        add_text(s, x + 0.14, 2.92, 1.8, 0.42, label, size=9.8, color=SMOKE)
    add_card(s, 0.95, 4.3, 5.55, 2.0, fill=WHITE, line=SANDAL)
    add_text(s, 1.2, 4.56, 2.3, 0.22, "SUSTAINABILITY", size=11, color=PEACOCK, bold=True)
    add_bullets(s, 1.2, 4.92, 4.8, 1.0, [
        "ASHA workers introduce Saheli during routine care visits.",
        "SHG leaders run community onboarding in small trusted groups.",
        "OSCs and legal-aid teams use it as a support layer, not a replacement.",
    ], size=11.2, color=SMOKE, bullet=PEACOCK)
    add_card(s, 6.82, 4.3, 5.5, 2.0, fill=RGBColor(35, 29, 47), line=RGBColor(104, 90, 118))
    add_text(s, 7.1, 4.56, 2.8, 0.22, "WHAT WE NEED TO IMPROVE FAST", size=11, color=HALDI, bold=True)
    add_bullets(s, 7.1, 4.92, 4.7, 1.0, [
        "Indic speech data from women in low-resource dialects.",
        "Verified district-level directories for OSCs, PHCs, legal aid, and Jan Aushadhi.",
        "Validation data for women’s-health screening pathways.",
    ], size=11.0, color=MIST, bullet=HALDI)


def slide_close(prs: Presentation):
    s = new_slide(prs)
    add_bg(s, "bg_cover.png")
    add_text(s, 0.85, 0.75, 2.2, 0.24, "WHY NOW", size=11, color=MIST, bold=True)
    add_text(s, 0.85, 1.2, 6.0, 1.0, "Build for the women current safety apps were never designed to hold.", font_name=FONT_HEAD, size=28, color=WHITE, bold=True)
    add_text(s, 0.85, 3.0, 5.4, 0.9, "Shared phones. Hindi-first interaction. Intermittent internet. A real need for proof, rights, and dignified support.", size=18, color=MIST)
    add_card(s, 0.85, 5.15, 5.1, 1.02, fill=RGBColor(39, 31, 53), line=RGBColor(105, 90, 121))
    add_text(s, 1.15, 5.48, 4.5, 0.34, "One conversation. Every right. Every scheme. Every tool.", size=18, color=WHITE, bold=True)
    add_text(s, 10.1, 6.35, 2.0, 0.34, "THANK YOU", font_name=FONT_HEAD, size=22, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)


def build():
    random.seed(11)
    build_backgrounds()
    prs = prs_new()
    slide_cover(prs)
    slide_insight(prs)
    slide_failure_modes(prs)
    slide_what_is(prs)
    slide_capabilities(prs)
    slide_stealth_flow(prs)
    slide_personas(prs)
    slide_stack(prs)
    slide_impact(prs)
    slide_close(prs)
    prs.save(str(OUT_PPTX))


if __name__ == "__main__":
    build()
